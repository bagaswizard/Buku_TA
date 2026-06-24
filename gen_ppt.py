#!/usr/bin/env python3
"""Generate 34-slide PPT for TA Sidang."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ============================================================
# CONSTANTS
# ============================================================
ITS_BLUE   = RGBColor(0x0A, 0x2F, 0x80)
LIGHT_BLUE = RGBColor(0x4A, 0x90, 0xD9)
BG_BLUE    = RGBColor(0xE8, 0xF0, 0xFE)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
DARK       = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT     = RGBColor(0x21, 0x96, 0xF3)
GRAY       = RGBColor(0x66, 0x66, 0x66)

BASE = os.path.dirname(os.path.abspath(__file__))
GAMBAR = os.path.join(BASE, 'gambar')
LOGO_ITS  = os.path.join(GAMBAR, 'Institut Teknologi Sepuluh Nopember - Blue.png')
LOGO_ELEK = os.path.join(GAMBAR, 'Elektro-trans.png')
LOGO_ELKA = os.path.join(GAMBAR, 'ELKA b202.png')

W = Inches(13.333)   # 16:9 width
H = Inches(7.5)      # 16:9 height

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

# ============================================================
# HELPER FUNCTIONS
# ============================================================

def add_logo_bar(slide):
    """Add 3 logos at top-right corner of every slide."""
    x_right = W - Inches(0.4)
    y_top   = Inches(0.15)
    h_logo  = Inches(0.4)
    for logo_path in [LOGO_ITS, LOGO_ELEK, LOGO_ELKA]:
        if os.path.exists(logo_path):
            x_right -= Inches(0.85)
            slide.shapes.add_picture(logo_path, x_right, y_top, height=h_logo)


def add_header_bar(slide, title_text):
    """Add a colored header bar with title text."""
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), W, Inches(0.9)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = ITS_BLUE
    bar.line.fill.background()
    tf = bar.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(26)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.LEFT
    tf.margin_left = Inches(0.6)
    tf.margin_top  = Inches(0.1)


def add_body_text(slide, bullets, left=0.8, top=1.3, width=11.7, font_size=16):
    """Add bullet-point body text."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(5.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = b
        p.font.size = Pt(font_size)
        p.font.color.rgb = DARK
        p.space_after = Pt(10)
        p.level = 0
    return tf


def add_subtitle_text(slide, text, top=2.8, font_size=18):
    """Centered subtitle text."""
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(top), Inches(12.3), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = GRAY
    p.alignment = PP_ALIGN.CENTER


def add_image_slide(slide, img_name, left=1.0, top=1.4, width=11.3, height=5.3):
    """Add an image to a slide, centering it in the available space."""
    path = os.path.join(GAMBAR, img_name)
    if os.path.exists(path):
        slide.shapes.add_picture(path, Inches(left), Inches(top),
                                 Inches(width), Inches(height))


def add_table_slide(slide, headers, rows, top=1.4, left=0.3):
    """Add a simple formatted table."""
    n_rows = len(rows) + 1
    n_cols = len(headers)
    tbl_width = Inches(12.7)
    tbl_height = Inches(min(0.5 * n_rows, 5.5))
    shape = slide.shapes.add_table(n_rows, n_cols, Inches(left), Inches(top),
                                   tbl_width, tbl_height)
    table = shape.table

    # Set column widths
    col_w = tbl_width // n_cols
    for i in range(n_cols):
        table.columns[i].width = col_w

    # Header row
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        for para in cell.text_frame.paragraphs:
            para.font.size = Pt(12)
            para.font.bold = True
            para.font.color.rgb = WHITE
            para.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = ITS_BLUE

    # Data rows
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = str(val)
            for para in cell.text_frame.paragraphs:
                para.font.size = Pt(11)
                para.font.color.rgb = DARK
                para.alignment = PP_ALIGN.CENTER
            if i % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = BG_BLUE


def new_slide():
    """Create blank slide"""
    layout = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(layout)


def title_slide(title, subtitle="", author=""):
    """Title slide with ITS blue background."""
    sl = new_slide()
    bg = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = ITS_BLUE
    bg.line.fill.background()

    # Title
    txBox = sl.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(2.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.LEFT

    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(20)
        p2.font.color.rgb = RGBColor(0xCC, 0xDD, 0xFF)
        p2.space_before = Pt(16)

    if author:
        txBox2 = sl.shapes.add_textbox(Inches(0.8), Inches(5.5), Inches(11.7), Inches(1.5))
        tf2 = txBox2.text_frame
        p3 = tf2.paragraphs[0]
        p3.text = author
        p3.font.size = Pt(18)
        p3.font.color.rgb = RGBColor(0xAA, 0xBB, 0xEE)

    # Logo at bottom-right
    if os.path.exists(LOGO_ITS):
        sl.shapes.add_picture(LOGO_ITS, Inches(11.0), Inches(6.0), height=Inches(1.0))
    return sl


def content_slide(title, bullets, img=None, img_size=(10.0, 4.5)):
    """Standard content slide with header, bullets, optional image."""
    sl = new_slide()
    add_header_bar(sl, title)
    add_logo_bar(sl)
    if img:
        add_body_text(sl, bullets, width=5.5, font_size=14)
        add_image_slide(sl, img, left=6.8, top=1.3, width=img_size[0], height=img_size[1])
    else:
        add_body_text(sl, bullets)
    return sl


def image_only_slide(title, img_name, img_width=11.3, img_height=5.3):
    """Slide with header and a large centered image."""
    sl = new_slide()
    add_header_bar(sl, title)
    add_logo_bar(sl)
    add_image_slide(sl, img_name, left=1.0, top=1.2, width=img_width, height=img_height)
    return sl


def blank_slide():
    sl = new_slide()
    add_logo_bar(sl)
    return sl


# ============================================================
# SLIDE 1 — JUDUL
# ============================================================
title_slide(
    "Sistem Navigasi Robot Otonom pada\nLingkungan Pasca Gempa Bertingkat\nmenggunakan Layered Costmap",
    "Tugas Akhir — Sidang",
    "Bagas | NRP xxxxx | Dosen Pembimbing: xxxxx\nDepartemen Teknik Elektro — Elektronika\nInstitut Teknologi Sepuluh Nopember"
)

# ============================================================
# BAB 1 — PENDAHULUAN (Slides 2-5)
# ============================================================

content_slide("Latar Belakang", [
    "Indonesia memiliki frekuensi gempa bumi tinggi",
    "Gempa menyebabkan perubahan struktural signifikan pada bangunan",
    "Lingkungan indoor pasca gempa tidak aman bagi manusia",
    "Robot otonom dapat berperan mengeksplorasi area berbahaya",
    "Robot memerlukan navigasi yang handal di lingkungan bertingkat pasca gempa",
    "Tantangan: known map tidak lagi akurat karena unknown obstacle",
])

content_slide("Latar Belakang — Kesenjangan Peta", [
    "Peta pra-gempa (known map) tersedia, tetapi...",
    "Pasca gempa: reruntuhan, pergeseran struktur, unknown obstacle",
    "Robot harus mampu mendeteksi dan menghindari unknown obstacle",
    "Lingkungan bertingkat tidak bisa direpresentasikan dalam 1 peta 2D",
    "Solusi: multi-region map + layered costmap + dynamic obstacle detection",
], img="sim_world.png", img_size=(6.5, 3.8))

content_slide("Rumusan Masalah", [
    "1. Bagaimana merepresentasikan lingkungan 3D bertingkat menjadi peta 2D?",
    "2. Bagaimana robot bernavigasi antar region yang terpisah?",
    "3. Bagaimana robot mendeteksi dan menghindari unknown obstacle pasca gempa?",
    "4. Bagaimana performa ICP pada berbagai kondisi lingkungan (normal vs gempa)?",
    "5. Bagaimana pengaruh Cost Scaling Factor terhadap kualitas global path?",
])

content_slide("Tujuan & Batasan", [
    "Mengembangkan sistem navigasi otonom untuk lingkungan pasca gempa bertingkat",
    "Mengimplementasikan layered costmap dengan 3 lapisan (global, local, transisi)",
    "Mengintegrasikan ICP, global planner, DWA planner, dan region switcher",
    "",
    "Batasan:",
    "   — Lingkungan simulasi MuJoCo, 3 region (A→B→C)",
    "   — Sensor LiDAR 2D, map dalam format occupancy grid",
    "   — Tidak ada fungsi Search and Rescue (SAR)",
    "   — Navigasi berbasis known map sebelum gempa",
])

# ============================================================
# BAB 2 — TINJAUAN PUSTAKA (Slides 6-9)
# ============================================================

content_slide("Occupancy Grid & Costmap", [
    "Occupancy Grid: representasi diskrit lingkungan dalam bentuk sel",
    "   — FREE_SPACE (0), LETHAL_OBSTACLE (254), NO_INFORMATION (255)",
    "   — Nilai 1–252: gradient cost hasil inflasi",
    "",
    "Layered Costmap: tiap layer menyimpan informasi berbeda",
    "   — Static Layer: obstacle dari known map",
    "   — Obstacle Layer: obstacle dari sensor LiDAR (dinamis)",
    "   — Inflation Layer: memperlebar obstacle untuk safety margin",
    "   — Transition Layer: informasi koneksi antar region",
])

content_slide("ICP — Iterative Closest Point", [
    "Algoritma untuk mencocokkan scan LiDAR dengan peta",
    "Prekomputasi KNN: O(1) lookup dengan hash table",
    "   — Setiap sel bebas di dekat obstacle disimpan K nearest occupied cells",
    "   — Dilakukan sekali saat startup, bukan saat runtime",
    "",
    "Proses ICP per iterasi:",
    "   1. Sample sensor points dengan random stride",
    "   2. Transform ke map frame  →  cari correspondences via KNN",
    "   3. Weighted SVD solver  →  hitung transformasi optimal",
    "   4. Akumulasi  →  cek konvergensi ($t_{norm}$, $r_{norm}$)",
])

content_slide("Path Planning & DWA Planner", [
    "Global Planner (Wavefront BFS):",
    "   — Propagasi potensial dari sel goal ke seluruh map",
    "   — Gradient descent dari posisi robot menuju goal",
    "   — Potensial kuadratik untuk aproksimasi Euclidean yang lebih akurat",
    "",
    "DWA (Dynamic Window Approach) Planner:",
    "   — Sample kecepatan dalam dynamic window",
    "   — Simulasikan trajectory dengan model kinematik",
    "   — Skor trajectory: obstacle cost, goal cost, path alignment, osilasi",
    "   — Pilih trajectory terbaik → kirim ke motor",
])

content_slide("Multi-Map Navigation", [
    "Lingkungan 3D bertingkat dibagi menjadi beberapa region 2D",
    "Setiap region adalah proyeksi 2D dari sub-bagian lingkungan 3D",
    "",
    "Key components:",
    "   — Combined Map: gabungan semua region dalam satu kanvas",
    "   — Transition Map: menyimpan wormhole connection antar region",
    "   — Region Switcher: memindahkan pose robot saat mencapai zona transisi",
    "   — Path Router: merute ulang goal ke zona transisi antar region",
])

# ============================================================
# BAB 3 — DESAIN & IMPLEMENTASI (Slides 10-22)
# ============================================================

content_slide("Arsitektur Sistem", [
    "Sensor: LiDAR → data scan untuk local costmap & ICP",
    "Map Server: combined map + transition map",
    "Costmap (3 lapis): Global (static+inflation), Local (obstacle), Transition",
    "Lokalisasi: ICP + prekomputasi KNN",
    "Planner: Global Planner → Path Router → DWA Planner",
    "Multi-Map: Region Switcher + Path Router untuk navigasi antar region",
], img="draft_navigation_flow.drawio.png", img_size=(6.6, 4.5))

image_only_slide("Alur Navigasi", "main_nav_flow.png", 10.0, 5.5)

content_slide("Multi-Region Map — Combined Map", [
    "Proses pembuatan combined map:",
    "   1. Scanning direktori  →  temukan semua file .pgm region",
    "   2. Deteksi ROI (Region of Interest) via bounding box",
    "   3. Ekstrak ROI dengan padding",
    "   4. Hitung grid layout (cols × rows)",
    "   5. Buat kanvas 2048×2048, paste ROI dengan centered alignment",
    "   6. Simpan combined_map.pgm",
], img="combined_map.png", img_size=(5.5, 4.0))

content_slide("Transition Map & Wormhole", [
    "Transition Map: menyimpan informasi koneksi antar region",
    "   — Pixel value 20: entry/exit point wormhole",
    "   — Pixel value 0:  jalur penghubung dalam region",
    "   — Pixel value 255: background (belum anotasi)",
    "",
    "Wormhole Connection: koneksi virtual antar dua region",
    "   — Koneksi bersifat linear (A→B→C, tidak bisa A→C)",
    "   — Algoritma Bresenham untuk menggambar jalur anotasi",
], img="transition_map_overlay.png", img_size=(5.5, 4.0))

content_slide("Prekomputasi KNN", [
    "Tujuan: mempercepat lookup correspondences ICP di runtime",
    "Dilakukan sekali saat menerima map baru (bukan saat runtime)",
    "",
    "Proses per sel obstacle:",
    "   1. Cek Moore neighbourhood 3×3 — ada sel bebas di sekitar?",
    "   2. Iterasi semua sel dalam radius r = 10 pixel",
    "   3. Untuk setiap sel bebas: insert obstacle ke sorted KNN list (k = 8)",
    "",
    "Hasil: setiap sel bebas memiliki daftar K nearest obstacle",
    "   → Lookup KNN saat runtime: O(1) via hash table",
])

content_slide("ICP — Pre-kalkulasi Cache", [
    r"Cos/sin cache untuk N beam pada sudut $\theta_0, \ldots, \theta_{N-1}$:",
    r"$\mathbf{C} = \begin{bmatrix} \cos\theta_0 & \cos\theta_1 & \cdots & \cos\theta_{N-1} \\ \sin\theta_0 & \sin\theta_1 & \cdots & \sin\theta_{N-1} \end{bmatrix}$",
    "",
    "Cache dihitung ulang hanya jika parameter scan berubah",
    "",
    "Diberikan range readings $r_0, \ldots, r_{N-1}$ → sensor point cloud:",
    r"$\mathbf{S} = \begin{bmatrix} r_0\cos\theta_0 & r_1\cos\theta_1 & \cdots & r_{N-1}\cos\theta_{N-1} \\ r_0\sin\theta_0 & r_1\sin\theta_1 & \cdots & r_{N-1}\sin\theta_{N-1} \end{bmatrix}$",
    "",
    "Untuk setiap $s_j$: cek map bounds → konversi ke flat index → lookup KNN",
])

content_slide("ICP — Loop Utama", [
    "Input: scan (N points), Output: transform T (new → old)",
    "",
    "T = Identity",
    "for iteration = 1 to max_iter (50):",
    "    1. random_scan = sampler(scan)            # stride s = 2",
    "    2. data = T · random_scan                 # ke map frame",
    "    3. M = matches(data, KNN)                 # cari correspondences",
    "    4. w = get_weights(M)                     # bobot tiap match",
    "    5. T_update = point_to_map(M, w)          # SVD solver",
    "    6. T = T_update · T                       # akumulasi",
    "    7. if converged(T_update): break",
    "return T",
    "",
    "Konvergensi: $||t|| < 0.01$m, $|\Delta\theta| < 0.01$rad ($\approx 0.57^\circ$)",
])

content_slide("Costmap — Global Costmap", [
    "Static Layer: obstacle dari known map (nilai 254/0/255)",
    "   — Dibangun sekali saat menerima combined map",
    "",
    "Inflation Layer: memperlebar obstacle dengan fungsi eksponensial",
    r"$C(d) = 253 \cdot e^{-\alpha \cdot (d - r_{inscribed})}$     untuk $d \le R$",
    "   — $\alpha$ = cost scaling factor (CSF)",
    "   — Semakin tinggi CSF → semakin curam penurunan → path lebih dekat obstacle",
    "",
    "Kedua layer menggunakan raytracing Bresenham untuk update cell",
])

content_slide("Costmap — Local & Transition", [
    "Local Costmap:",
    "   — Obstacle Layer: dibangun dari scan LiDAR secara real-time",
    "   — Inflation Layer: sama seperti global, tapi dari data obstacle layer",
    "   — Rolling window mengikuti posisi robot",
    "   — Mampu mendeteksi unknown obstacle (tidak ada di known map)",
    "",
    "Transition Costmap:",
    "   — Dibangun dari transition map (hasil anotasi wormhole)",
    "   — Menyimpan informasi koneksi antar region",
    "   — Digunakan oleh Region Switcher untuk deteksi zona transisi",
])

content_slide("Raytracing — Bresenham", [
    "Digunakan untuk marking sel yang dilalui sinar LiDAR sebagai FREE",
    "Juga digunakan dalam Pixel Annotator untuk wormhole connections",
    "",
    "Inisialisasi:",
    r"   $\Delta x = |x_1 - x_0|,\quad \Delta y = |y_1 - y_0|$",
    r"   $s_x = \text{sign}(x_1 - x_0),\quad s_y = \text{sign}(y_1 - y_0)$",
    r"   $\text{err} = \Delta x - \Delta y$",
    "",
    "Iterasi sepanjang garis:",
    "   — Update error $e_2 = 2\cdot\text{err}$",
    "   — Langkah horizontal/vertikal/diagonal berdasarkan error",
    "   — Panggil fungsi raytrace untuk setiap pixel yang dilalui",
])

content_slide("Global Planner — Wavefront BFS", [
    "Wavefront BFS dari sel goal ke seluruh map:",
    "   — Setiap sel menyimpan jarak potensial ke goal",
    "   — Menggunakan potensial kuadratik untuk akurasi Euclidean",
    "",
    "Cost function:",
    r"   $c_{eff} = \text{costmap}[n] \cdot f + c_{neutral}$",
    "",
    "Gradient descent:",
    "   — Hitung gradien dari potential map",
    "   — Normalisasi → step size → iterasi sampai mencapai goal",
    "   — Menghasilkan urutan waypoints",
])

content_slide("DWA Planner — Trajectory Scoring", [
    "Dynamic Window: batas kecepatan yang dapat dicapai dalam 1 control cycle",
    r"   $v_{min} = \max(v_{limit}, v_{current} - \dot{v}_{max}\cdot\Delta t)$",
    "",
    "Trajectory simulation:",
    r"   $x_{t+1} = x_t + (v_x\cos\theta_t + v_y\cos(\pi/2+\theta_t))\cdot\Delta t$",
    r"   $y_{t+1} = y_t + (v_x\sin\theta_t + v_y\sin(\pi/2+\theta_t))\cdot\Delta t$",
    r"   $\theta_{t+1} = \theta_t + \omega\cdot\Delta t$",
    "",
    "Total cost:",
    r"   $C_{total} = C_{osc} + w_{obs}C_{obs} + w_{gf}C_{gf} + w_{align}C_{align} + w_{goal}C_{goal}$",
])

content_slide("Region Switcher & Path Router", [
    "Path Router (Path Interceptor):",
    "   — Menerima navigation goal akhir",
    "   — Jika goal di region berbeda → beri goal perantara ke zona transisi",
    "   — Meneruskan hasil ke Global Planner",
    "",
    "Region Switcher (Costmap Jumper):",
    "   — Melacak region robot saat ini (dari ICP + transition map)",
    "   — Saat robot di zona transisi: pindahkan pose ke region berikutnya",
    "   — Memperbarui posisi awal ICP pada region baru",
    "   — Memuat ulang costmap untuk region baru",
])

# ============================================================
# BAB 4 — PENGUJIAN & ANALISIS (Slides 23-31)
# ============================================================

content_slide("Skenario Pengujian", [
    "Lingkungan simulasi: MuJoCo — 3 region (A, B, C)",
    "Kondisi pengujian:",
    "   — Normal: lantai rata, sesuai known map",
    "   — Gempa: lantai tidak rata + unknown obstacle (simulasi reruntuhan)",
    "",
    "Metrik pengujian:",
    "   — ICP: error x, y, yaw per satuan waktu",
    "   — Path: panjang, clearance, waypoints, cross track error",
    "   — Tracking: mean & max cross track error (DWA)",
    "",
    "Cost Scaling Factor yang diuji: 10, 20, 50, 100",
], img="sim_RA.png", img_size=(6.0, 3.6))

image_only_slide("Hasil ICP — Region A (Normal)", "icp_errors_RA.png", 10.5, 4.8)

content_slide("Hasil ICP — Region A (Normal)", [
    "Error ICP pada sumbu x, y, dan θ (yaw)",
    "Secara umum: performa stabil dengan error rendah di sepanjang lintasan",
    "Lonjakan error signifikan di sekitar t ≈ 50 detik",
    "Penyebab: robot berada di zona transisi antar region",
    "   → Lingkungan yang dapat dideteksi LiDAR sangat terbatas",
    "   → Kualitas scan matching menurun",
    "   → Akumulasi error pada estimasi transformasi",
])

content_slide("Hasil ICP — Gempa & Region C", [
    "Region A — Lantai Tidak Rata (Gempa):",
    "   — Peningkatan osilasi signifikan pada sumbu x dan yaw",
    "   — Permukaan tidak rata → scan LiDAR bervariasi vertikal → scan matching menurun",
    "   — Lonjakan t≈55 tetap terjadi (zona transisi)",
    "",
    "Region C:",
    "   — Error max pada ketiga sumbu lebih tinggi dibanding Region A",
    "   — Lebih tidak stabil secara umum",
    "   — Penyebab: lebih banyak obstacle (KNN lookup lebih besar)",
    "   — Trajectory lebih kompleks dengan banyak rotasi",
])

content_slide("Hasil Path — Region A", [
    "CSF 10: path cenderung di tengah antar obstacle → lebih aman, lebih panjang",
    "CSF 100: path lebih dekat ke obstacle → lebih pendek, lebih lurus",
    "CSF 20 & 50: karakteristik peralihan",
    "",
    "CSF 100 menghasilkan path paling efisien untuk Region A",
    "   — Panjang path lebih pendek (2.538 m vs 2.596 m)",
    "   — Waypoints lebih sedikit (509 vs 521)",
    "   — Min clearance tetap memadai (0.216 m)",
], img="path_all_RA.png", img_size=(5.5, 3.8))

content_slide("Hasil Path — Region B", [
    "Region B: hanya berupa lorong sempit",
    "Semua CSF menghasilkan path yang hampir identik",
    "   → Tidak ada ruang terbuka bagi variasi cost untuk mengubah lintasan",
    "",
    "Tidak ada pengaruh signifikan CSF terhadap path di Region B",
    "Semua path identik — tidak ada deviasi yang berarti antar CSF",
], img="path_all_RB.png", img_size=(5.5, 3.8))

content_slide("Hasil Path — Region C", [
    "CSF 10: selalu di tengah, lebih aman, lebih panjang",
    "CSF 100: path lebih pendek, terutama di area terbuka",
    "CSF 20 & 50: di antara kedua ekstrem",
    "",
    "CSF 100 menghasilkan path paling cocok untuk Region C:",
    "   — Memanfaatkan ruang terbuka untuk jalur lebih pendek",
    "   — Clearance terhadap obstacle masih memadai",
    "   — Keseimbangan terbaik antara efisiensi dan keamanan",
], img="path_all_RC.png", img_size=(5.2, 3.6))

content_slide("Perbandingan vs Ground Truth", [
    "Ground truth = path referensi ideal (konservatif, aman)",
    "",
    "Region A: CSF 100 lebih mendekati ground truth (deviasi rendah)",
    "Region B: semua CSF menghasilkan deviasi hampir sama (lorong sempit)",
    "Region C: CSF 100 deviasi paling tinggi — trade-off panjang vs akurasi",
    "",
    "Kesimpulan: CSF tinggi → path lebih pendek, tapi tidak selalu lebih mirip ground truth",
    "   → Pemilihan CSF bergantung pada prioritas: keamanan vs efisiensi",
])

content_slide("Hasil Tracking — DWA Planner", [
    "DWA Planner mengikuti global path sambil menghindari obstacle",
    "",
    "Kondisi Normal:",
    "   — DWA mengikuti global path dengan baik",
    "   — Mean cross track error rendah (0.01–0.14 m)",
    "",
    "Kondisi Gempa (Region A):",
    "   — Max cross track error meningkat signifikan (3.13 m vs 0.03 m)",
    "   — Unknown obstacle di lintasan → DWA menghindar menjauhi global path",
    "   — Simpangan ini tercermin sebagai max cross track error yang tinggi",
])

# ============================================================
# BAB 5 — PENUTUP (Slides 32-34)
# ============================================================

content_slide("Kesimpulan", [
    "1. Sistem navigasi berhasil bernavigasi dari Region A ke Region C",
    "   melalui zona transisi menggunakan Region Switcher",
    "",
    "2. Layered costmap (global, local, transition) efektif untuk:",
    "   — Representasi multi-region",
    "   — Deteksi dan penghindaran unknown obstacle",
    "   — Transisi mulus antar region",
    "",
    "3. ICP bekerja akurat di sebagian besar lintasan, error meningkat di zona transisi",
    "   — Lantai tidak rata menambah osilasi pada sumbu x dan yaw",
    "",
    "4. Cost Scaling Factor 100 menghasilkan path paling efisien (pendek, lurus, clearance memadai)",
    "",
    "5. DWA Planner mampu mendeteksi dan menghindari unknown obstacle secara real-time",
])

content_slide("Saran", [
    "1. Penggabungan sensor IMU untuk meningkatkan akurasi yaw ICP",
    "",
    "2. Dynamic adjustment CSF berdasarkan tipe region (lorong vs ruang terbuka)",
    "",
    "3. Adaptive stride sampling ICP untuk mengurangi lonjakan error di zona transisi",
    "",
    "4. Multi-robot coordination untuk eksplorasi paralel di lingkungan nyata",
    "",
    "5. Implementasi dan pengujian pada robot fisik",
    "",
    "6. Integrasi dengan 3D SLAM untuk rekonstruksi lingkungan",
])

# --- Thank You Slide ---
sl = new_slide()
bg = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
bg.fill.solid()
bg.fill.fore_color.rgb = ITS_BLUE
bg.line.fill.background()

txBox = sl.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(11.7), Inches(2.5))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Terima Kasih"
p.font.size = Pt(52)
p.font.color.rgb = WHITE
p.font.bold = True
p.alignment = PP_ALIGN.CENTER

add_subtitle_text(sl, "Questions & Answers", top=4.5, font_size=24)

if os.path.exists(LOGO_ITS):
    sl.shapes.add_picture(LOGO_ITS, Inches(5.5), Inches(5.5), height=Inches(1.2))

# ============================================================
# SAVE
# ============================================================
output_path = os.path.join(BASE, 'SIDANG_TA.pptx')
prs.save(output_path)
print(f"PPT saved: {output_path}")
print(f"Total slides: {len(prs.slides)}")
