#!/usr/bin/env python3
"""Generate PowerPoint presentation for Sidang Tugas Akhir."""

import os
import copy
import csv

from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree

# =============================================================================
# DESIGN SYSTEM CONSTANTS
# =============================================================================
PRIMARY = RGBColor(0x00, 0x3A, 0x70)       # #003A70 - Biru ITS
SECONDARY = RGBColor(0x4A, 0x90, 0xD9)     # #4A90D9 - Biru muda
HEADING = RGBColor(0x1A, 0x2B, 0x4C)       # #1A2B4C - Judul
BODY = RGBColor(0x33, 0x33, 0x33)          # #333333 - Teks body
ORANGE = RGBColor(0xF0, 0x8C, 0x2E)        # #F08C2E - Aksen
BOX_FILL = RGBColor(0xFF, 0xF8, 0xE7)     # #FFF8E7 - Cream card bg
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF5, 0xF9, 0xFF)     # Kartu pertanyaan
WARN_BG = RGBColor(0xFF, 0xF8, 0xF0)      # Panel batasan
DARK_BG = RGBColor(0x00, 0x2B, 0x54)      # Slide 1 gradient fallback
SUBTLE_TEXT = RGBColor(0x55, 0x55, 0x55)  # Captions
BORDER_GRAY = RGBColor(0xE0, 0xE0, 0xE0)

GRAD_OUTLINE_S   = RGBColor(0xE2, 0xEE, 0xFC)
GRAD_OUTLINE_E   = RGBColor(0xBE, 0xD4, 0xEF)
GRAD_CONTEXT_S   = RGBColor(0xE4, 0xF0, 0xFF)
GRAD_CONTEXT_E   = RGBColor(0xC1, 0xD7, 0xF3)
GRAD_TECH_S      = RGBColor(0xED, 0xF4, 0xFC)
GRAD_TECH_E      = RGBColor(0xD2, 0xE3, 0xF3)
GRAD_RESULTS_S   = RGBColor(0xFB, 0xF3, 0xE8)
GRAD_RESULTS_E   = RGBColor(0xE4, 0xD2, 0xB8)
GRAD_CLOSING_S    = RGBColor(0xE8, 0xF1, 0xFC)
GRAD_CLOSING_E    = RGBColor(0xCD, 0xDF, 0xF1)

FONT_TITLE = "Montserrat"
FONT_BODY = "Lato"
FONT_CODE = "Fira Code"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

MARGIN = Inches(0.93)  # ~7% margin

# =============================================================================
# HELPER FUNCTIONS
# =============================================================================

def set_slide_bg(slide, color):
    """Set solid color background for a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def set_gradient_bg(slide, start_color, end_color, angle_deg=135):
    """Set gradient background for a slide using lxml for gradient stops."""
    cSld = slide._element.find(qn('p:cSld'))

    bg = cSld.find(qn('p:bg'))
    if bg is None:
        bg = etree.Element(qn('p:bg'))
        spTree = cSld.find(qn('p:spTree'))
        if spTree is not None:
            spTree.addprevious(bg)
        else:
            cSld.insert(0, bg)
    else:
        for child in list(bg):
            bg.remove(child)

    bgPr = etree.SubElement(bg, qn('p:bgPr'))

    gradFill = etree.SubElement(bgPr, qn('a:gradFill'))
    gradFill.set('rotWithShape', '1')

    gsLst = etree.SubElement(gradFill, qn('a:gsLst'))

    gs1 = etree.SubElement(gsLst, qn('a:gs'))
    gs1.set('pos', '0')
    srgb1 = etree.SubElement(gs1, qn('a:srgbClr'))
    srgb1.set('val', str(start_color))

    gs2 = etree.SubElement(gsLst, qn('a:gs'))
    gs2.set('pos', '100000')
    srgb2 = etree.SubElement(gs2, qn('a:srgbClr'))
    srgb2.set('val', str(end_color))

    lin = etree.SubElement(gradFill, qn('a:lin'))
    lin.set('ang', str(angle_deg * 60000))
    lin.set('scaled', '0')


def add_rect(slide, left, top, width, height, fill_color=None, border_color=None, border_width=None, corner_radius=None):
    """Add a rectangle shape."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if corner_radius else MSO_SHAPE.RECTANGLE,
        left, top, width, height
    )
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.fill.solid()
        if border_width:
            shape.line.width = border_width
    return shape


def add_text_box(slide, left, top, width, height, text, font_name=FONT_BODY,
                 font_size=Pt(16), color=BODY, bold=False, alignment=PP_ALIGN.LEFT,
                 line_spacing=1.2, anchor=MSO_ANCHOR.TOP, italic=False):
    """Add a text box with single paragraph."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.text_frame.word_wrap = True
    txBox.text_frame.auto_size = None
    p = txBox.text_frame.paragraphs[0]
    p.text = text
    p.font.name = font_name
    p.font.size = font_size
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.italic = italic
    p.alignment = alignment
    p.space_after = Pt(0)
    p.space_before = Pt(0)
    p.line_spacing = Pt(font_size.pt * line_spacing)
    return txBox


def add_rich_text_box(slide, left, top, width, height, paragraphs_data,
                      anchor=MSO_ANCHOR.TOP):
    """Add a text box with multiple paragraphs. paragraphs_data is a list of dicts:
    {text, font_name, font_size, color, bold, alignment, line_spacing, space_after, italic}
    """
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.text_frame.word_wrap = True
    txBox.text_frame.auto_size = None

    for i, pdata in enumerate(paragraphs_data):
        if i == 0:
            p = txBox.text_frame.paragraphs[0]
        else:
            p = txBox.text_frame.add_paragraph()

        p.text = pdata.get("text", "")
        p.font.name = pdata.get("font_name", FONT_BODY)
        p.font.size = pdata.get("font_size", Pt(16))
        p.font.color.rgb = pdata.get("color", BODY)
        p.font.bold = pdata.get("bold", False)
        p.font.italic = pdata.get("italic", False)
        p.alignment = pdata.get("alignment", PP_ALIGN.LEFT)
        ls = pdata.get("line_spacing", 1.3)
        if isinstance(ls, (int, float)):
            p.line_spacing = Pt(pdata["font_size"].pt * ls)
        else:
            p.line_spacing = ls
        p.space_after = pdata.get("space_after", Pt(4))

    return txBox


def add_slide_number(slide, num):
    """Add slide number at bottom right."""
    add_text_box(slide, SLIDE_W - MARGIN - Inches(0.5), SLIDE_H - Inches(0.45),
                 Inches(0.5), Inches(0.3), str(num),
                 font_size=Pt(10), color=BODY,
                 alignment=PP_ALIGN.RIGHT)


def add_divider(slide, left, top, width, color=ORANGE, height=Pt(2)):
    """Add a horizontal divider line."""
    shape = add_rect(slide, left, top, width, height, fill_color=color)


def add_card(slide, left, top, width, height, number, title, description,
             card_bg=BOX_FILL, accent=SECONDARY):
    """Add a card with number circle, title, and description."""
    card = add_rect(slide, left, top, width, height, fill_color=card_bg, corner_radius=Cm(0.2))
    # Number circle
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, left + Inches(0.2), top + Inches(0.2),
        Inches(0.5), Inches(0.5)
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = accent
    circle.line.fill.background()
    tf = circle.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = str(number)
    p.font.name = FONT_TITLE
    p.font.size = Pt(16)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # Title
    add_text_box(slide, left + Inches(0.85), top + Inches(0.15),
                 width - Inches(1.1), Inches(0.4),
                 title, font_name=FONT_TITLE, font_size=Pt(14),
                 color=HEADING, bold=True)

    # Description
    add_text_box(slide, left + Inches(0.25), top + Inches(0.7),
                 width - Inches(0.5), height - Inches(0.85),
                 description, font_size=Pt(12), color=BODY,
                 line_spacing=1.35)


def add_icon_bullet(slide, left, top, width, height, icon_char, title, body_text,
                    accent=SECONDARY):
    """Add an icon-style bullet point."""
    add_rect(slide, left, top, width, height,
             fill_color=BOX_FILL, corner_radius=Cm(0.10))
    # Icon circle
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, left, top, Inches(0.45), Inches(0.45)
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = accent
    circle.line.fill.background()
    tf = circle.text_frame
    p = tf.paragraphs[0]
    p.text = icon_char
    p.font.size = Pt(14)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    add_text_box(slide, left + Inches(0.6), top - Inches(0.02),
                 width - Inches(0.6), Inches(0.25),
                 title, font_name=FONT_TITLE, font_size=Pt(14),
                 color=HEADING, bold=True)

    add_text_box(slide, left + Inches(0.6), top + Inches(0.25),
                 width - Inches(0.6), height - Inches(0.25),
                 body_text, font_size=Pt(12), color=BODY,
                 line_spacing=1.3)


# =============================================================================
# SLIDE FUNCTIONS
# =============================================================================


def _add_csv_table(slide, left, top, col_widths, csv_rel_path, title=None):
    """Read CSV from lampiran/ directory and add a table shape to the slide.
    Returns (table_shape, total_height_inches)."""
    base_dir = os.path.dirname(os.path.abspath(__file__))
    csv_path = os.path.join(base_dir, "lampiran", csv_rel_path)

    with open(csv_path, newline='') as f:
        reader = csv.DictReader(f)
        headers = reader.fieldnames
        rows = list(reader)

    n_rows = 1 + len(rows)   # header + data
    n_cols = len(headers)

    row_h = Inches(0.28)
    tbl_w = sum(col_widths)
    tbl_h = n_rows * row_h

    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, tbl_w, tbl_h)
    table = table_shape.table

    # Set column widths
    for ci, w in enumerate(col_widths):
        table.columns[ci].width = w

    # Header row
    for ci, h in enumerate(headers):
        cell = table.cell(0, ci)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.name = FONT_TITLE
            p.font.size = Pt(8)
            p.font.color.rgb = WHITE
            p.font.bold = True
            p.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = PRIMARY

    # Data rows
    for ri, row in enumerate(rows):
        for ci, h in enumerate(headers):
            cell = table.cell(ri + 1, ci)
            cell.text = row[h]
            for p in cell.text_frame.paragraphs:
                p.font.name = FONT_BODY
                p.font.size = Pt(8)
                p.font.color.rgb = BODY
                p.alignment = PP_ALIGN.CENTER
            if ri % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = BOX_FILL

    # Title above table
    actual_top = top
    if title:
        add_text_box(slide, left, top - Inches(0.25), tbl_w, Inches(0.22),
                     title, font_name=FONT_TITLE, font_size=Pt(9),
                     color=HEADING, bold=True)
        actual_top = top

    return table_shape, tbl_h + (Inches(0.25) if title else 0), actual_top


def slide_title(prs):
    """Slide 1: Title slide — overhang logo bar with clean layered layout."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    base_dir = os.path.dirname(os.path.abspath(__file__))

    # =========================================================================
    # Background: two large circles for depth
    # =========================================================================
    bg_circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(-2.5), Inches(-3.0), Inches(9.0), Inches(9.0)
    )
    bg_circle.fill.solid()
    bg_circle.fill.fore_color.rgb = RGBColor(0x0A, 0x3C, 0x62)
    bg_circle.line.fill.background()

    bg_circle2 = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, SLIDE_W - Inches(3.5), SLIDE_H - Inches(4.5),
        Inches(6.0), Inches(6.0)
    )
    bg_circle2.fill.solid()
    bg_circle2.fill.fore_color.rgb = RGBColor(0x0A, 0x3C, 0x62)
    bg_circle2.line.fill.background()

    # =========================================================================
    # Logo bar — shifted partially off-screen so only bottom-left rounding shows
    # =========================================================================
    LOGO_BAR_W = Inches(2.65)
    LOGO_BAR_H = Inches(0.68)
    LOGO_RADIUS = Cm(0.22)

    OVERHANG_Y = Inches(0.18)
    OVERHANG_X = Inches(0.75)
    logo_bar_x = SLIDE_W - LOGO_BAR_W + OVERHANG_X
    logo_bar_y = Inches(0.0) - OVERHANG_Y

    # Drop shadow — offset behind, matching overhang
    SHADOW_OFFSET = Pt(4)
    shadow = add_rect(slide,
                      logo_bar_x - SHADOW_OFFSET,
                      logo_bar_y + SHADOW_OFFSET,
                      LOGO_BAR_W + SHADOW_OFFSET,
                      LOGO_BAR_H + Pt(2),
                      fill_color=RGBColor(0x00, 0x1A, 0x3A),
                      corner_radius=LOGO_RADIUS)
    shadow.line.fill.background()

    # White foreground bar
    logo_bg = add_rect(slide, logo_bar_x, logo_bar_y, LOGO_BAR_W, LOGO_BAR_H,
                       fill_color=WHITE, corner_radius=LOGO_RADIUS)
    logo_bg.line.fill.background()

    # =========================================================================
    # Logos — three images, uniform reduced height, inside the white bar
    # =========================================================================
    LOGO_H = Inches(0.40)
    LOGO_GAP = Inches(0.12)
    LOGO_PAD_LEFT = Inches(0.12)
    LOGO_CENTER_Y = Inches(0.10)

    LOGO_Y = LOGO_CENTER_Y - Pt(3)

    its_path = os.path.join(base_dir, "gambar", "Institut Teknologi Sepuluh Nopember - Blue.png")
    its_w = LOGO_H * (1400.0 / 902.0)
    its_x = logo_bar_x + LOGO_PAD_LEFT
    slide.shapes.add_picture(its_path, its_x, LOGO_Y, its_w, LOGO_H)

    elektro_path = os.path.join(base_dir, "gambar", "Elektro-trans.png")
    el_w = LOGO_H * (503.0 / 496.0)
    el_x = its_x + its_w + LOGO_GAP
    slide.shapes.add_picture(elektro_path, el_x, LOGO_Y, el_w, LOGO_H)

    elka_path = os.path.join(base_dir, "gambar", "ELKA b202.png")
    elka_w = LOGO_H * (701.0 / 588.0)
    elka_x = el_x + el_w + LOGO_GAP
    slide.shapes.add_picture(elka_path, elka_x, LOGO_Y, elka_w, LOGO_H)

    # =========================================================================
    # Decorative accent — orange vertical bar aligned with title
    # =========================================================================
    add_rect(slide, Inches(0.55), Inches(1.25), Pt(5), Inches(2.1), fill_color=ORANGE)

    # =========================================================================
    # Title block — positioned below logo bar with clear gap
    # =========================================================================
    TITLE_SIZE = Pt(30)
    title_x = Inches(1.05)
    title_y = Inches(1.25)

    add_rich_text_box(slide, title_x, title_y, Inches(9.0), Inches(2.5), [
        {"text": "NAVIGASI ROBOT BERBASIS", "font_name": FONT_TITLE, "font_size": TITLE_SIZE,
         "color": WHITE, "bold": True, "line_spacing": 1.22},
        {"text": "LAYERED COSTMAP", "font_name": FONT_TITLE, "font_size": TITLE_SIZE,
         "color": WHITE, "bold": True, "line_spacing": 1.22},
        {"text": "UNTUK PERGERAKAN ANTAR LANTAI", "font_name": FONT_TITLE, "font_size": TITLE_SIZE,
         "color": RGBColor(0xC8, 0xDC, 0xEE), "bold": True, "line_spacing": 1.22},
        {"text": "PADA LINGKUNGAN INDOOR PASCA GEMPA BUMI", "font_name": FONT_TITLE, "font_size": TITLE_SIZE,
         "color": RGBColor(0xA0, 0xBE, 0xDA), "bold": True, "line_spacing": 1.3, "space_after": Pt(28)},
    ])

    # =========================================================================
    # Divider — inside title block area, below text
    # =========================================================================
    add_rect(slide, title_x, Inches(3.50), Inches(3.6), Pt(2.5), fill_color=ORANGE)

    # =========================================================================
    # Author info — positioned below divider with clean gap
    # =========================================================================
    add_rich_text_box(slide, title_x, Inches(3.95), Inches(7.0), Inches(2.0), [
        {"text": "Bagas Surya Wirawan", "font_name": FONT_TITLE, "font_size": Pt(22),
         "color": WHITE, "bold": True, "line_spacing": 1.15},
        {"text": "NRP. 5022221026", "font_name": FONT_BODY, "font_size": Pt(16),
         "color": RGBColor(0xBB, 0xD5, 0xEC), "line_spacing": 1.0},
        {"text": "", "font_name": FONT_BODY, "font_size": Pt(8), "color": WHITE, "line_spacing": 1.0},
        {"text": "Dosen Pembimbing:", "font_name": FONT_BODY, "font_size": Pt(14),
         "color": RGBColor(0xCC, 0xCC, 0xCC), "bold": True, "line_spacing": 1.05},
        {"text": "Dr. Ir. Djoko Purwanto, M.Eng.", "font_name": FONT_BODY, "font_size": Pt(14),
         "color": RGBColor(0xBB, 0xD5, 0xEC), "line_spacing": 0.95},
        {"text": "Fajar Budiman, S.T., M.Sc.", "font_name": FONT_BODY, "font_size": Pt(14),
         "color": RGBColor(0xBB, 0xD5, 0xEC), "line_spacing": 0.95},
        {"text": "", "font_name": FONT_BODY, "font_size": Pt(8), "color": WHITE, "line_spacing": 1.0},
        {"text": "Departemen Teknik Elektro", "font_name": FONT_BODY, "font_size": Pt(12),
         "color": RGBColor(0x9B, 0xBB, 0xDB), "line_spacing": 0.95},
        {"text": "Fakultas Teknologi Elektro dan Informatika Cerdas", "font_name": FONT_BODY, "font_size": Pt(12),
         "color": RGBColor(0x9B, 0xBB, 0xDB), "line_spacing": 0.95},
        {"text": "Institut Teknologi Sepuluh Nopember  ·  Surabaya  ·  2026", "font_name": FONT_BODY, "font_size": Pt(12),
         "color": RGBColor(0x85, 0xA8, 0xCA), "line_spacing": 0.95},
    ])

    # =========================================================================
    # Bottom accent
    # =========================================================================
    add_rect(slide, Inches(0.55), SLIDE_H - Inches(0.35), Inches(2.5), Pt(1.5),
             fill_color=RGBColor(0x14, 0x4D, 0x77))


def slide_outline(prs):
    """Slide 2: Outline / Daftar Isi."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_gradient_bg(slide, GRAD_OUTLINE_S, GRAD_OUTLINE_E)

    add_text_box(slide, MARGIN, Inches(0.4), Inches(4), Inches(0.6),
                 "DAFTAR ISI", font_name=FONT_TITLE, font_size=Pt(30),
                 color=HEADING, bold=True)
    add_divider(slide, MARGIN, Inches(0.95), Inches(1.8))

    items = [
        ("01", "Latar Belakang", "Konteks gempa bumi di Indonesia dan urgensi robot SAR otonom"),
        ("02", "Rumusan Masalah & Tujuan", "Pertanyaan penelitian, batasan, dan tujuan yang ingin dicapai"),
        ("03", "Metodologi & Desain Sistem", "Arsitektur navigasi: costmap, ICP, A*, dan DWA planner"),
        ("04", "Hasil Pengujian", "Evaluasi ICP, costmap, dan navigasi penuh pada arena simulasi"),
        ("05", "Kesimpulan & Saran", "Temuan utama dan arah pengembangan selanjutnya"),
    ]

    card_w = Inches(3.8)
    card_h = Inches(0.95)
    start_x = MARGIN
    start_y = Inches(1.35)
    gap_x = Inches(0.25)
    gap_y = Inches(0.18)

    positions = [
        (start_x, start_y),
        (start_x, start_y + card_h + gap_y),
        (start_x + card_w + gap_x, start_y),
        (start_x + card_w + gap_x, start_y + card_h + gap_y),
        (start_x + 2 * (card_w + gap_x), start_y),
    ]

    for (x, y), (num, title, desc) in zip(positions, items):
        card = add_rect(slide, x, y, card_w, card_h, fill_color=BOX_FILL, corner_radius=Cm(0.2))
        # Big number
        add_text_box(slide, x + Inches(0.2), y + Inches(0.1),
                     Inches(0.7), Inches(0.6),
                     num, font_name=FONT_TITLE, font_size=Pt(28),
                     color=SECONDARY, bold=True)
        add_text_box(slide, x + Inches(0.85), y + Inches(0.1),
                     card_w - Inches(1.1), Inches(0.35),
                     title, font_name=FONT_TITLE, font_size=Pt(15),
                     color=HEADING, bold=True)
        add_text_box(slide, x + Inches(0.85), y + Inches(0.48),
                     card_w - Inches(1.1), Inches(0.4),
                     desc, font_size=Pt(11), color=SUBTLE_TEXT,
                     line_spacing=1.25)

    add_slide_number(slide, 2)


def slide_background(prs):
    """Slide 3: Latar Belakang."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_gradient_bg(slide, GRAD_CONTEXT_S, GRAD_CONTEXT_E)

    add_text_box(slide, MARGIN, Inches(0.35), Inches(5), Inches(0.55),
                 "LATAR BELAKANG", font_name=FONT_TITLE, font_size=Pt(30),
                 color=HEADING, bold=True)
    add_divider(slide, MARGIN, Inches(0.85), Inches(2.2))

    # Left side: 5 focused bullet points — full width
    bullet_data = [
        ("1", "INDONESIA RAWAN GEMPA BUMI",
         "Terletak di 4 lempeng tektonik (Pasifik, Eurasia, Indo-Australia, Laut Filipina). "
         "Rata-rata 18 gempa/hari (BMKG, Sabtaji 2020)."),
        ("2", "OPERASI SAR & PERAN ROBOT",
         "Petugas SAR hadapi ancaman keselamatan. Robot darat digunakan "
         "pasca gempa Jepang untuk inspeksi (Lin et al., 2022). Robot UAV memberi "
         "gambaran medan. Robot otonom krusial saat infrastruktur komunikasi rusak."),
        ("3", "TANTANGAN NAVIGASI INDOOR",
         "Puing, lantai tidak rata, tangga. Robot perlu lokalisasi, "
         "path planning, dan mencapai tujuan di lingkungan tidak terstruktur. "
         "Navigasi indoor masih menjadi permasalahan robot otonom."),
        ("4", "SIMULASI ARENA KRSRI 2024",
         "Sulit diuji di kondisi nyata → simulasi Mujoco. Arena dari "
         "Kontes Robot SAR Indonesia (KRSRI) 2024. Platform: robot hexapod + LiDAR 2D. "
         "Keterbatasan LiDAR 2D pada lingkungan bertingkat."),
        ("5", "SOLUSI: LAYERED COSTMAP",
         "Informasi 3D → occupancy grid + beberapa costmap 2D. "
         "Setiap layer = satu lantai/region. Layer transisi menghubungkan antar lantai. "
         "Robot berpindah layer sesuai posisi di bangunan."),
    ]

    y = Inches(1.1)
    for num, title, desc in bullet_data:
        draw_background_bullet(slide, MARGIN, y, Inches(11.5), Inches(1.0), num, title, desc)
        y += Inches(1.08)

    # Bottom highlight — one key takeaway
    add_rect(slide, MARGIN, Inches(6.55), Inches(11.5), Inches(0.45),
             fill_color=RGBColor(0xFF, 0xF3, 0xE0), corner_radius=Cm(0.1))
    add_rect(slide, MARGIN, Inches(6.55), Pt(4), Inches(0.45), fill_color=ORANGE)
    add_text_box(slide, MARGIN + Inches(0.12), Inches(6.55), Inches(11.3), Inches(0.45),
                 "Layered costmap menyederhanakan kompleksitas 3D menjadi representasi 2D "
                 "yang manageable dengan sensor LiDAR 2D pada robot hexapod.",
                 font_size=Pt(11), color=BODY, italic=True, line_spacing=1.2)

    add_slide_number(slide, 3)


def draw_background_bullet(slide, left, top, width, height, num, title, desc):
    """Draw a numbered bullet point for latar belakang."""
    add_rect(slide, left, top, width, height,
             fill_color=BOX_FILL, corner_radius=Cm(0.10))
    # Number circle
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, left, top + Inches(0.05), Inches(0.4), Inches(0.4)
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = SECONDARY
    circle.line.fill.background()
    tf = circle.text_frame
    p = tf.paragraphs[0]
    p.text = num
    p.font.name = FONT_TITLE
    p.font.size = Pt(14)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    add_text_box(slide, left + Inches(0.55), top, width - Inches(0.55), Inches(0.25),
                 title, font_name=FONT_TITLE, font_size=Pt(13), color=HEADING, bold=True)
    add_text_box(slide, left + Inches(0.55), top + Inches(0.27), width - Inches(0.55),
                 height - Inches(0.27),
                 desc, font_size=Pt(11), color=BODY, line_spacing=1.3)


def slide_problem(prs):
    """Slide 4: Rumusan Masalah."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_gradient_bg(slide, GRAD_CONTEXT_S, GRAD_CONTEXT_E)

    add_text_box(slide, MARGIN, Inches(0.35), Inches(5), Inches(0.55),
                 "RUMUSAN MASALAH", font_name=FONT_TITLE, font_size=Pt(30),
                 color=HEADING, bold=True)
    add_divider(slide, MARGIN, Inches(0.85), Inches(2.5))

    # Three question cards
    questions = [
        ("01", "REPRESENTASI LINGKUNGAN 3D",
         "Bagaimana cara menggunakan layered costmap untuk mewakilkan informasi "
         "lingkungan 3D yang terdiri dari beberapa lantai sehingga bisa digunakan "
         "pada algoritma navigasi 2D?"),
        ("02", "ALGORITMA NAVIGASI",
         "Bagaimana merancang algoritma lokalisasi, path planning, dan path tracking "
         "yang bekerja pada sistem navigasi berbasis layered costmap?"),
        ("03", "VALIDASI SISTEM",
         "Apakah robot darat otonom dapat melakukan navigasi hingga mencapai tujuan "
         "pada lingkungan simulasi menggunakan metode navigasi antar lantai yang dikembangkan?"),
    ]

    card_w = Inches(11.5)
    card_h = Inches(1.35)
    card_x = MARGIN
    card_y = Inches(1.15)

    for i, (num, title, desc) in enumerate(questions):
        y = card_y + i * (card_h + Inches(0.25))
        card = add_rect(slide, card_x, y, card_w, card_h,
                        fill_color=LIGHT_BG, corner_radius=Cm(0.15))
        # Left accent bar
        add_rect(slide, card_x, y, Pt(5), card_h, fill_color=SECONDARY)
        # Number
        add_text_box(slide, card_x + Inches(0.2), y + Inches(0.15),
                     Inches(0.6), Inches(0.55),
                     num, font_name=FONT_TITLE, font_size=Pt(28),
                     color=SECONDARY, bold=True)
        # Title
        add_text_box(slide, card_x + Inches(0.8), y + Inches(0.12),
                     Inches(4.0), Inches(0.4),
                     title, font_name=FONT_TITLE, font_size=Pt(15),
                     color=HEADING, bold=True)
        # Description
        add_text_box(slide, card_x + Inches(0.8), y + Inches(0.55),
                     card_w - Inches(1.1), Inches(0.7),
                     desc, font_size=Pt(13), color=BODY, line_spacing=1.35)

    add_slide_number(slide, 4)


def slide_constraints(prs):
    """Slide 5: Batasan Penelitian."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_gradient_bg(slide, GRAD_CONTEXT_S, GRAD_CONTEXT_E)

    add_text_box(slide, MARGIN, Inches(0.35), Inches(5), Inches(0.55),
                 "BATASAN PENELITIAN", font_name=FONT_TITLE, font_size=Pt(30),
                 color=HEADING, bold=True)
    add_divider(slide, MARGIN, Inches(0.85), Inches(2.5))

    constraints = [
        ("Simulasi Arena KRSRI 2024 (Mujoco)",
         "Merekayasa kondisi pasca gempa dengan rintangan tidak diketahui,\n"
         "tanjakan, dan permukaan tidak rata pada SKALA LEBIH KECIL."
         "\n\nBukan kondisi gempa sebenarnya → validasi awal metode.",
         PRIMARY),
        ("Platform Robot Hexapod + LiDAR 2D",
         "Robot hexapod 6 kaki bernavigasi secara OTONOM penuh dari\n"
         "titik awal ke tujuan yang telah ditentukan."
         "\n\nKeterbatasan LiDAR 2D pada lingkungan bertingkat\n"
         "menjadi tantangan dalam sistem navigasi.",
         SECONDARY),
        ("Lingkungan Statis Diketahui (Known Map)",
         "Algoritma navigasi bekerja pada kondisi lingkungan DENGAN\n"
         "rintangan statis yang SUDAH DIKETAHUI."
         "\n\nRintangan dinamis tidak termasuk dalam cakupan\n"
         "pengujian tugas akhir ini.",
         ORANGE),
    ]

    card_w = Inches(11.5)
    card_h = Inches(1.55)
    card_x = MARGIN
    card_y = Inches(1.2)

    for i, (title, desc, color) in enumerate(constraints):
        y = card_y + i * (card_h + Inches(0.25))
        add_rect(slide, card_x, y, card_w, card_h,
                 fill_color=BOX_FILL, corner_radius=Cm(0.15))
        add_rect(slide, card_x, y, Pt(5), card_h, fill_color=color)

        # Warning triangle
        tri = slide.shapes.add_shape(
            MSO_SHAPE.ISOSCELES_TRIANGLE,
            card_x + Inches(0.25), y + Inches(0.2),
            Inches(0.35), Inches(0.32)
        )
        tri.fill.solid()
        tri.fill.fore_color.rgb = color
        tri.line.fill.background()
        # "!" in triangle
        tf = tri.text_frame
        p = tf.paragraphs[0]
        p.text = "!"
        p.font.name = FONT_TITLE
        p.font.size = Pt(14)
        p.font.color.rgb = WHITE
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

        # Title
        add_text_box(slide, card_x + Inches(0.8), y + Inches(0.12),
                     Inches(5.0), Inches(0.35),
                     title, font_name=FONT_TITLE, font_size=Pt(16),
                     color=color, bold=True)
        # Description
        add_text_box(slide, card_x + Inches(0.8), y + Inches(0.48),
                     card_w - Inches(1.1), Inches(1.0),
                     desc, font_size=Pt(12), color=BODY, line_spacing=1.35)

    # Bottom note
    add_rect(slide, MARGIN, Inches(6.55), Inches(11.5), Inches(0.45),
             fill_color=RGBColor(0xFF, 0xF3, 0xE0), corner_radius=Cm(0.1))
    add_rect(slide, MARGIN, Inches(6.55), Pt(4), Inches(0.45), fill_color=ORANGE)
    add_text_box(slide, MARGIN + Inches(0.12), Inches(6.55), Inches(11.3), Inches(0.45),
                 "Metode yang dikembangkan TIDAK dirancang untuk langsung digunakan di kondisi pasca gempa. "
                 "Pengujian dilakukan pada simulasi untuk validasi awal sebelum penerapan di skenario nyata.",
                 font_size=Pt(11), color=BODY, italic=True, line_spacing=1.2)

    add_slide_number(slide, 5)


def slide_objectives(prs):
    """Slide 6: Tujuan & Manfaat."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_gradient_bg(slide, GRAD_CONTEXT_S, GRAD_CONTEXT_E)

    add_text_box(slide, MARGIN, Inches(0.35), Inches(5), Inches(0.55),
                 "TUJUAN & MANFAAT", font_name=FONT_TITLE, font_size=Pt(30),
                 color=HEADING, bold=True)
    add_divider(slide, MARGIN, Inches(0.85), Inches(2.2))

    # Left column: Tujuan
    add_text_box(slide, MARGIN, Inches(1.1), Inches(2.5), Inches(0.4),
                 "Tujuan Penelitian", font_name=FONT_TITLE, font_size=Pt(18),
                 color=PRIMARY, bold=True)

    goals = [
        ("🎯", "Representasi Multi-Lantai",
         "Mengembangkan metode representasi lingkungan menggunakan layered costmap untuk "
         "menyederhanakan informasi lingkungan multi-lantai"),
        ("🎯", "Algoritma Navigasi Antar Lantai",
         "Merancang dan mengimplementasikan algoritma lokalisasi, path planning, dan path tracking "
         "yang menggunakan layered costmap"),
        ("🎯", "Evaluasi di Simulasi",
         "Melakukan pengujian sistem di lingkungan simulasi untuk mengevaluasi efektivitas "
         "dan keterbatasan pendekatan yang diusulkan"),
    ]

    y = Inches(1.6)
    for icon, title, desc in goals:
        add_icon_goal(slide, MARGIN + Inches(0.1), y, Inches(5.8), Inches(1.1), title, desc)
        y += Inches(1.25)

    # Right column: Manfaat
    right_x = Inches(7.0)
    add_text_box(slide, right_x, Inches(1.1), Inches(2.5), Inches(0.4),
                 "Manfaat Penelitian", font_name=FONT_TITLE, font_size=Pt(18),
                 color=PRIMARY, bold=True)

    benefits = [
        ("⭐", "Kontribusi pada sistem navigasi robot otonom di lingkungan kompleks"),
        ("⭐", "Dasar awal navigasi robot di area bertingkat atau runtuh pasca gempa"),
        ("⭐", "Fondasi penelitian lanjutan untuk robot SAR otonom di area bencana"),
    ]

    by = Inches(1.6)
    for icon, desc in benefits:
        add_icon_bullet(slide, right_x + Inches(0.1), by, Inches(5.8), Inches(0.55),
                        icon, "", desc)
        by += Inches(0.95)

    add_slide_number(slide, 6)


def add_icon_goal(slide, left, top, width, height, title, desc):
    """Draw a goal item."""
    add_rect(slide, left, top, width, height,
             fill_color=BOX_FILL, corner_radius=Cm(0.10))
    # Icon circle
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, left, top + Inches(0.02), Inches(0.4), Inches(0.4)
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = SECONDARY
    circle.line.fill.background()
    tf = circle.text_frame
    p = tf.paragraphs[0]
    p.text = "•"
    p.font.size = Pt(16)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    add_text_box(slide, left + Inches(0.55), top, width - Inches(0.55), Inches(0.25),
                 title, font_name=FONT_TITLE, font_size=Pt(13), color=HEADING, bold=True)
    add_text_box(slide, left + Inches(0.55), top + Inches(0.28), width - Inches(0.55),
                 height - Inches(0.28),
                 desc, font_size=Pt(11), color=BODY, line_spacing=1.35)


def slide_literature(prs):
    """Slide 7: Dasar Teori."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_gradient_bg(slide, GRAD_TECH_S, GRAD_TECH_E)

    add_text_box(slide, MARGIN, Inches(0.35), Inches(5), Inches(0.55),
                 "DASAR TEORI", font_name=FONT_TITLE, font_size=Pt(30),
                 color=HEADING, bold=True)
    add_divider(slide, MARGIN, Inches(0.85), Inches(2.2))

    # 5 subsection titles from bab/2-tinjauan-pustaka.tex
    topics = [
        "Occupancy Grid",
        "Layered Costmap",
        "Lokalisasi Iterative Closest Point (ICP)",
        "Path Planning A*",
        "Path Tracking",
    ]

    card_w = Inches(5.5)
    card_h = Inches(1.1)
    gap_x = Inches(0.4)
    gap_y = Inches(0.25)
    start_x = MARGIN
    start_y = Inches(1.2)
    cols = 2

    colors = [PRIMARY, PRIMARY, SECONDARY, SECONDARY, ORANGE]

    for i, topic in enumerate(topics):
        row = i // cols
        col = i % cols
        x = start_x + col * (card_w + gap_x)
        y = start_y + row * (card_h + gap_y)

        add_rect(slide, x, y, card_w, card_h,
                 fill_color=BOX_FILL, corner_radius=Cm(0.15))
        add_rect(slide, x, y, Pt(5), card_h, fill_color=colors[i])

        # Number
        add_text_box(slide, x + Inches(0.15), y + Inches(0.12),
                     Inches(0.6), Inches(0.45),
                     f"{i+1:02d}", font_name=FONT_TITLE, font_size=Pt(24),
                     color=colors[i], bold=True)
        # Title
        add_text_box(slide, x + Inches(0.8), y + Inches(0.12),
                     card_w - Inches(1.0), Inches(0.85),
                     topic, font_name=FONT_TITLE, font_size=Pt(16),
                     color=HEADING, bold=True)

    add_slide_number(slide, 7)


def slide_prior_research(prs):
    """Slide 8: Penelitian Terdahulu — Comparison Table."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_gradient_bg(slide, GRAD_TECH_S, GRAD_TECH_E)

    add_text_box(slide, MARGIN, Inches(0.35), Inches(8), Inches(0.55),
                 "PENELITIAN TERDAHULU", font_name=FONT_TITLE, font_size=Pt(28),
                 color=HEADING, bold=True)
    add_divider(slide, MARGIN, Inches(0.85), Inches(2.5))

    rows = 6
    cols = 2
    table_left = MARGIN
    table_top = Inches(1.15)
    table_w = Inches(11.5)
    row_h = Inches(0.95)

    table_shape = slide.shapes.add_table(rows, cols, table_left, table_top,
                                          table_w, row_h * rows)
    table = table_shape.table

    table.columns[0].width = Inches(5.4)
    table.columns[1].width = Inches(6.1)

    def set_cell(cell, text, font_size=Pt(9), color=BODY, bold=False,
                 fill_color=None, alignment=PP_ALIGN.LEFT):
        cell.text = ""
        p = cell.text_frame.paragraphs[0]
        p.text = text
        p.font.name = FONT_BODY
        p.font.size = font_size
        p.font.color.rgb = color
        p.font.bold = bold
        p.alignment = alignment
        cell.text_frame.word_wrap = True
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_left = Pt(5)
        cell.margin_right = Pt(5)
        cell.margin_top = Pt(3)
        cell.margin_bottom = Pt(3)
        if fill_color:
            cell.fill.solid()
            cell.fill.fore_color.rgb = fill_color

    # Header
    set_cell(table.cell(0, 0), "Pendekatan Unik pada Penelitian Terdahulu",
             font_size=Pt(11), color=WHITE, bold=True, fill_color=PRIMARY,
             alignment=PP_ALIGN.CENTER)
    set_cell(table.cell(0, 1), "Perbandingan dengan Penelitian Ini",
             font_size=Pt(11), color=WHITE, bold=True, fill_color=PRIMARY,
             alignment=PP_ALIGN.CENTER)

    # Data rows
    data = [
        ("Macenski et al. (2023)\nLayered Costmap",
         "Mengusulkan arsitektur layered costmap pertama — membagi costmap "
         "menjadi layer semantik terpisah (static, obstacle, inflation) yang "
         "memperbarui master grid secara terstruktur.",
         "Mengadopsi layered costmap dan memperluas dengan tiga layer spesifik: "
         "Global Costmap, Local Costmap, serta Transition Costmap yang menyimpan "
         "informasi koneksi dan zona transisi antar region."),

        ("Hu et al. (2022)\nAdaptive Slope Navigation",
         "Multi-layer costmap adaptif untuk navigasi permukaan berkemiringan. "
         "Algoritma deteksi slope mengubah peta biaya agar kemiringan dianggap "
         "dapat dilalui (dengan biaya tambahan), bukan sebagai obstacle.",
         "Fokus berbeda: Hu menangani variasi kemiringan, sedangkan penelitian ini "
         "menangani transisi antar lantai. Keduanya menyederhanakan informasi 3D "
         "menjadi representasi costmap 2D bertingkat."),

        ("Kim et al. (2024)\nIntegrated Navigation Map",
         "Membuat peta navigasi terintegrasi dengan menggabungkan point cloud "
         "multi-lantai dan topologi node graph. Mengandalkan elevator sebagai "
         "penghubung antar lantai pada gedung bertingkat.",
         "Tidak bergantung pada elevator — menggunakan mekanisme region "
         "transition dan zona transisi yang terdeteksi dari combined map. "
         "Lebih adaptif untuk skenario pasca gempa."),

        ("Palac\u00edn et al. (2023)\nInter-floor ICP + Elevator",
         "Navigasi antar lantai menggunakan elevator dengan lokalisasi berbasis "
         "ICP. Evaluasi performa ICP dilakukan pada kondisi pintu lift terbuka "
         "maupun tertutup.",
         "Menggunakan ICP yang sama dengan optimasi prekomputasi KNN "
         "(hash table, lookup O(1), k=8, r=10). Dilengkapi region switcher "
         "untuk mendeteksi perpindahan region secara otomatis dari pose robot."),

        ("Jung et al. (2024)\nA* dengan Realistic Cost",
         "A* dengan fungsi biaya realistis — mempertimbangkan faktor seperti "
         "waktu tunggu lift. Robot dapat memilih antara lift atau tangga "
         "berdasarkan perhitungan biaya total.",
         "A* dilengkapi Path Interceptor yang memecah jalur multi-region menjadi "
         "segmen per region. Cost Scaling Factor (CSF=100) menghasilkan jalur "
         "terpendek dengan tetap menjaga obstacle clearance optimal."),
    ]

    alt_fill = RGBColor(0xF2, 0xF7, 0xFC)

    for i, (label, left_text, right_text) in enumerate(data):
        row_fill = alt_fill if i % 2 == 0 else WHITE

        set_cell(table.cell(i + 1, 0), f"{label}\n\n{left_text}",
                 font_size=Pt(9), color=BODY, fill_color=row_fill,
                 alignment=PP_ALIGN.LEFT)
        set_cell(table.cell(i + 1, 1), right_text,
                 font_size=Pt(9), color=BODY, fill_color=row_fill,
                 alignment=PP_ALIGN.LEFT)

        # Bold the author line in column 0
        p = table.cell(i + 1, 0).text_frame.paragraphs[0]
        # Split and format first line as bold/colored
        cell0 = table.cell(i + 1, 0)
        cell0.text = ""
        tf = cell0.text_frame
        p1 = tf.paragraphs[0]
        p1.text = label
        p1.font.name = FONT_TITLE
        p1.font.size = Pt(10)
        p1.font.color.rgb = PRIMARY
        p1.font.bold = True
        p1.alignment = PP_ALIGN.LEFT
        p1.space_after = Pt(2)
        p2 = tf.add_paragraph()
        p2.text = left_text
        p2.font.name = FONT_BODY
        p2.font.size = Pt(9)
        p2.font.color.rgb = BODY
        p2.font.bold = False
        p2.alignment = PP_ALIGN.LEFT

    add_slide_number(slide, 8)


def slide_system_overview(prs):
    """Slide 9: Arsitektur Sistem."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_gradient_bg(slide, GRAD_TECH_S, GRAD_TECH_E)

    add_text_box(slide, MARGIN, Inches(0.35), Inches(5), Inches(0.55),
                 "ARSITEKTUR SISTEM", font_name=FONT_TITLE, font_size=Pt(30),
                 color=HEADING, bold=True)
    add_divider(slide, MARGIN, Inches(0.85), Inches(2.2))

    base_dir = os.path.dirname(os.path.abspath(__file__))
    arch_path = os.path.join(base_dir, "gambar", "navigation_architecture.png")

    img_h = Inches(5.6)
    img_w = int(img_h * (1177.0 / 774.0))
    img_left = (SLIDE_W - img_w) / 2
    img_top = Inches(1.15)

    # Rectangle frame — sharp, bold outline around the image
    add_rect(slide, img_left, img_top, img_w, img_h,
             fill_color=None, border_color=PRIMARY, border_width=Pt(3))

    # The architecture image
    slide.shapes.add_picture(arch_path, img_left, img_top, img_w, img_h)

    add_slide_number(slide, 9)


def slide_nav_flow(prs):
    """Slide 10: Detail Alur Navigasi."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_gradient_bg(slide, GRAD_TECH_S, GRAD_TECH_E)

    add_text_box(slide, MARGIN, Inches(0.35), Inches(7), Inches(0.55),
                 "DETAIL ALUR NAVIGASI", font_name=FONT_TITLE, font_size=Pt(26),
                 color=HEADING, bold=True)
    add_divider(slide, MARGIN, Inches(0.85), Inches(2.5))

    base_dir = os.path.dirname(os.path.abspath(__file__))
    flow_path = os.path.join(base_dir, "gambar", "main_nav_flow_split.png")
    flow_h = Inches(5.5)
    flow_w = int(flow_h * (752.0 / 542.0))
    flow_left = (SLIDE_W - flow_w) / 2
    flow_top = Inches(1.2)

    add_rect(slide, flow_left, flow_top, flow_w, flow_h,
             fill_color=None, border_color=PRIMARY, border_width=Pt(3))
    slide.shapes.add_picture(flow_path, flow_left, flow_top, flow_w, flow_h)

    add_slide_number(slide, 10)


def slide_map_processing(prs):
    """Slide 11: Pengolahan Data Map (Multi-Region)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_gradient_bg(slide, GRAD_TECH_S, GRAD_TECH_E)

    add_text_box(slide, MARGIN, Inches(0.35), Inches(7), Inches(0.55),
                 "PENGOLAHAN DATA MAP", font_name=FONT_TITLE, font_size=Pt(28),
                 color=HEADING, bold=True)
    add_divider(slide, MARGIN, Inches(0.85), Inches(2.0))

    base_dir = os.path.dirname(os.path.abspath(__file__))

    # --- Top: Combined Map (full arena) ---
    combined_path = os.path.join(base_dir, "gambar", "combined_map.png")
    combined_h = Inches(2.9)
    combined_w = combined_h  # 1:1
    combined_left = (SLIDE_W - combined_w) / 2
    combined_top = Inches(1.25)

    add_text_box(slide, combined_left, combined_top - Inches(0.28),
                 combined_w, Inches(0.25),
                 "Combined Map — Gabungan Region A, B, C",
                 font_name=FONT_TITLE, font_size=Pt(11), color=PRIMARY, bold=True,
                 alignment=PP_ALIGN.CENTER)
    slide.shapes.add_picture(combined_path, combined_left, combined_top,
                              combined_w, combined_h)

    # --- Bottom: Transition Map & Overlay side by side ---
    trans_h = Inches(2.35)
    trans_path = os.path.join(base_dir, "gambar", "transition_map.png")
    trans_w = trans_h  # 512x512 = 1:1
    overlay_path = os.path.join(base_dir, "gambar", "transition_map_overlay.png")
    overlay_w = int(trans_h * (728.0 / 733.0))

    gap = Inches(0.6)
    total_w = trans_w + gap + overlay_w
    trans_left = (SLIDE_W - total_w) / 2
    overlay_left = trans_left + trans_w + gap
    trans_top = combined_top + combined_h + Inches(0.28)

    # Labels
    add_text_box(slide, trans_left, trans_top - Inches(0.28),
                 trans_w, Inches(0.25),
                 "Transition Map", font_name=FONT_TITLE, font_size=Pt(11),
                 color=PRIMARY, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, overlay_left, trans_top - Inches(0.28),
                 overlay_w, Inches(0.25),
                 "Transition Map Overlay", font_name=FONT_TITLE, font_size=Pt(11),
                 color=ORANGE, bold=True, alignment=PP_ALIGN.CENTER)

    # Images
    slide.shapes.add_picture(trans_path, trans_left, trans_top, trans_w, trans_h)
    slide.shapes.add_picture(overlay_path, overlay_left, trans_top, overlay_w, trans_h)

    # Caption
    add_text_box(slide, MARGIN, Inches(6.82), Inches(11.5), Inches(0.20),
                 "Transition Map menyimpan koneksi spasial (A↔B, B↔C), orientasi region, dan pose awal robot",
                 font_size=Pt(9), color=BODY, alignment=PP_ALIGN.CENTER)

    add_slide_number(slide, 11)


def slide_costmap_overview(prs):
    """Slide 12: Costmap — Gambaran Umum."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_gradient_bg(slide, GRAD_TECH_S, GRAD_TECH_E)

    add_text_box(slide, MARGIN, Inches(0.35), Inches(7), Inches(0.55),
                 "COSTMAP \u2014 REPRESENTASI BIAYA NAVIGASI", font_name=FONT_TITLE, font_size=Pt(26),
                 color=HEADING, bold=True)
    add_divider(slide, MARGIN, Inches(0.85), Inches(2.5))

    # Left: definition + layer composition table
    add_text_box(slide, MARGIN, Inches(1.1), Inches(6.5), Inches(0.8),
                 "Costmap adalah grid 2D yang menyimpan informasi biaya navigasi. "
                 "Nilai rendah menandakan area aman, nilai tinggi menandakan area "
                 "berbahaya atau terlarang. Sistem ini membagi costmap menjadi "
                 "tiga lapisan dengan fungsi berbeda:",
                 font_size=Pt(11), color=BODY, line_spacing=1.35)

    # Layer composition table as cards
    layers_info = [
        ("GLOBAL COSTMAP", "Static Layer + Inflation Layer", PRIMARY,
         "Perencanaan path skala luas dari peta tetap"),
        ("LOCAL COSTMAP", "Obstacle Layer + Inflation Layer", SECONDARY,
         "Navigasi jangka pendek, real-time dari LiDAR"),
        ("TRANSITION COSTMAP", "Static Layer + Transition Layer", ORANGE,
         "Perpindahan antar region, zona transisi"),
    ]

    ly = Inches(2.0)
    for title, comp, color, func in layers_info:
        add_rect(slide, MARGIN, ly, Inches(6.5), Inches(0.75),
                 fill_color=BOX_FILL, corner_radius=Cm(0.10))
        add_rect(slide, MARGIN, ly, Pt(5), Inches(0.75), fill_color=color)
        add_text_box(slide, MARGIN + Inches(0.2), ly + Inches(0.05),
                     Inches(2.5), Inches(0.3),
                     title, font_name=FONT_TITLE, font_size=Pt(12),
                     color=color, bold=True)
        add_text_box(slide, MARGIN + Inches(2.8), ly + Inches(0.05),
                     Inches(3.5), Inches(0.3),
                     comp, font_size=Pt(10), color=HEADING, bold=True)
        add_text_box(slide, MARGIN + Inches(0.2), ly + Inches(0.38),
                     Inches(6.1), Inches(0.3),
                     func, font_size=Pt(9), color=SUBTLE_TEXT)
        ly += Inches(0.85)

    # Right side: costmap constants table
    right_x = Inches(7.5)
    add_rect(slide, right_x, Inches(1.1), Inches(5.3), Inches(5.2),
             fill_color=BOX_FILL, corner_radius=Cm(0.15))
    add_text_box(slide, right_x + Inches(0.2), Inches(1.15), Inches(4.9), Inches(0.3),
                 "Nilai Konstanta Costmap", font_name=FONT_TITLE,
                 font_size=Pt(14), color=HEADING, bold=True)

    consts = [
        ("NO_INFORMATION", "255", "Area tidak diketahui"),
        ("LETHAL_OBSTACLE", "254", "Obstacle mematikan"),
        ("INSCRIBED_INFLATED", "253", "Batas inflasi maksimal"),
        ("INFLATED_OBSTACLE", "1\u2013252", "Area inflasi obstacle"),
        ("TRANSITION_CELL", "1\u20135", "Zona transisi region"),
        ("FREE_SPACE", "0", "Ruang bebas"),
    ]

    vy = Inches(1.6)
    for name, val, desc in consts:
        add_text_box(slide, right_x + Inches(0.2), vy, Inches(2.0), Inches(0.22),
                     name, font_size=Pt(8), color=HEADING, bold=True)
        add_text_box(slide, right_x + Inches(2.3), vy, Inches(0.8), Inches(0.22),
                     val, font_size=Pt(9), color=PRIMARY, bold=True,
                     alignment=PP_ALIGN.CENTER)
        add_text_box(slide, right_x + Inches(3.2), vy, Inches(1.8), Inches(0.22),
                     desc, font_size=Pt(8), color=SUBTLE_TEXT)
        vy += Inches(0.28)

    # Bottom note
    add_rect(slide, right_x + Inches(0.2), vy + Inches(0.10), Inches(4.9), Inches(0.65),
             fill_color=RGBColor(0xFD, 0xF2, 0xE3), corner_radius=Cm(0.10))
    add_text_box(slide, right_x + Inches(0.35), vy + Inches(0.15), Inches(4.6), Inches(0.55),
                 "Inscribed radius adalah jarak dari pusat robot ke batas "
                 "footprint. Sel dalam radius ini dianggap terlalu dekat "
                 "dan diberi cost maksimal 253.",
                 font_size=Pt(8), color=SUBTLE_TEXT, line_spacing=1.3)

    add_slide_number(slide, 12)


def slide_costmap_global(prs):
    """Slide 13: Global Costmap."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_gradient_bg(slide, GRAD_TECH_S, GRAD_TECH_E)

    add_text_box(slide, MARGIN, Inches(0.35), Inches(7), Inches(0.55),
                 "GLOBAL COSTMAP \u2014 PERENCANAAN JALUR SKALA LUAS",
                 font_name=FONT_TITLE, font_size=Pt(26),
                 color=HEADING, bold=True)
    add_divider(slide, MARGIN, Inches(0.85), Inches(2.5))

    # Left content
    add_text_box(slide, MARGIN, Inches(1.1), Inches(6.5), Inches(0.45),
                 "Global costmap mencakup seluruh area yang diketahui dari peta. "
                 "Digunakan oleh global planner (A*) untuk merencanakan path "
                 "dari posisi robot menuju tujuan dalam skala luas. "
                 "Costmap ini tidak berubah secara dinamis selama navigasi.",
                 font_size=Pt(11), color=BODY, line_spacing=1.35)

    # Static layer card
    add_rect(slide, MARGIN, Inches(1.7), Inches(6.5), Inches(1.1),
             fill_color=BOX_FILL, corner_radius=Cm(0.15))
    add_rect(slide, MARGIN, Inches(1.7), Pt(5), Inches(1.1), fill_color=PRIMARY)
    add_text_box(slide, MARGIN + Inches(0.2), Inches(1.75), Inches(3.0), Inches(0.3),
                 "Static Layer", font_name=FONT_TITLE, font_size=Pt(14),
                 color=PRIMARY, bold=True)
    add_text_box(slide, MARGIN + Inches(0.2), Inches(2.05), Inches(6.1), Inches(0.6),
                 "Menyimpan obstacle tetap dari peta. Mengkonversi nilai occupancy grid "
                 "(0\u2013100) ke cost internal (0\u2013255): unknown=255, free=0, "
                 "occupied=254. Bersumber dari combined map dan tidak berubah.",
                 font_size=Pt(10), color=BODY, line_spacing=1.3)

    # Inflation layer card
    add_rect(slide, MARGIN, Inches(2.95), Inches(6.5), Inches(1.1),
             fill_color=BOX_FILL, corner_radius=Cm(0.15))
    add_rect(slide, MARGIN, Inches(2.95), Pt(5), Inches(1.1), fill_color=PRIMARY)
    add_text_box(slide, MARGIN + Inches(0.2), Inches(3.0), Inches(3.0), Inches(0.3),
                 "Inflation Layer", font_name=FONT_TITLE, font_size=Pt(14),
                 color=PRIMARY, bold=True)
    add_text_box(slide, MARGIN + Inches(0.2), Inches(3.30), Inches(6.1), Inches(0.6),
                 "Memperluas obstacle dengan gradient cost untuk menjaga jarak aman. "
                 "Cost Scaling Factor (CSF) memengaruhi kecuraman gradien: "
                 "CSF tinggi = penurunan cost lebih curam, path lebih mepet ke obstacle. "
                 "CSF rendah = gradien landai, path menjauhi obstacle.",
                 font_size=Pt(10), color=BODY, line_spacing=1.3)

    # Right side: inflation formula
    right_x = Inches(7.5)
    add_rect(slide, right_x, Inches(1.1), Inches(5.3), Inches(3.0),
             fill_color=RGBColor(0xFD, 0xF2, 0xE3), corner_radius=Cm(0.15))
    add_text_box(slide, right_x + Inches(0.2), Inches(1.15), Inches(4.9), Inches(0.3),
                 "Fungsi Inflasi", font_name=FONT_TITLE, font_size=Pt(14),
                 color=ORANGE, bold=True)
    add_text_box(slide, right_x + Inches(0.2), Inches(1.55), Inches(4.9), Inches(1.5),
                 "C(d) = 253 \u00b7 e^(\u2212\u03b1 \u00b7 (d \u2212 r_inscribed))\n"
                 "untuk r_inscribed < d \u2264 R\n\n"
                 "d = jarak dari obstacle\n"
                 "R = inflation radius\n"
                 "r_inscribed = inscribed radius (cost = 253)\n"
                 "\u03b1 = scaling factor (kontrol kecuraman)\n\n"
                 "C(d) = 253 untuk d \u2264 r_inscribed\n"
                 "C(d) = 0 untuk d > R",
                 font_size=Pt(11), color=BODY, line_spacing=1.35)

    # CSF note
    add_rect(slide, right_x, Inches(4.35), Inches(5.3), Inches(0.5),
             fill_color=BOX_FILL, corner_radius=Cm(0.10))
    add_text_box(slide, right_x + Inches(0.15), Inches(4.38), Inches(5.0), Inches(0.45),
                 "Cost Scaling Factor (CSF) mengalikan costmap sebelum inflasi. "
                 "Nilai uji: 10, 20, 50, 100. CSF = 100 menghasilkan path terpendek.",
                 font_size=Pt(9), color=BODY, line_spacing=1.3)

    add_slide_number(slide, 13)


def slide_costmap_local(prs):
    """Slide 14: Local Costmap."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_gradient_bg(slide, GRAD_TECH_S, GRAD_TECH_E)

    add_text_box(slide, MARGIN, Inches(0.35), Inches(7), Inches(0.55),
                 "LOCAL COSTMAP \u2014 NAVIGASI JANGKA PENDEK REAL-TIME",
                 font_name=FONT_TITLE, font_size=Pt(26),
                 color=HEADING, bold=True)
    add_divider(slide, MARGIN, Inches(0.85), Inches(2.5))

    # Left content
    add_text_box(slide, MARGIN, Inches(1.1), Inches(6.5), Inches(0.45),
                 "Local costmap mencakup area di sekitar robot dalam radius "
                 "tertentu dan bergerak mengikuti robot (rolling window). "
                 "Digunakan oleh DWA local planner untuk perencanaan gerak "
                 "jangka pendek dan menghindari rintangan real-time.",
                 font_size=Pt(11), color=BODY, line_spacing=1.35)

    # Rolling window card
    add_rect(slide, MARGIN, Inches(1.7), Inches(6.5), Inches(0.85),
             fill_color=BOX_FILL, corner_radius=Cm(0.15))
    add_rect(slide, MARGIN, Inches(1.7), Pt(5), Inches(0.85), fill_color=SECONDARY)
    add_text_box(slide, MARGIN + Inches(0.15), Inches(1.73), Inches(3.0), Inches(0.25),
                 "Rolling Window", font_name=FONT_TITLE, font_size=Pt(13),
                 color=SECONDARY, bold=True)
    add_text_box(slide, MARGIN + Inches(0.15), Inches(1.98), Inches(6.2), Inches(0.5),
                 "Grid bergerak mengikuti posisi robot. Origin digeser setiap "
                 "kali robot berpindah sehingga robot selalu di dekat pusat. "
                 "Sel baru yang muncul di area belum tercakup diisi dari data LiDAR.",
                 font_size=Pt(10), color=BODY, line_spacing=1.25)

    # Obstacle layer card
    add_rect(slide, MARGIN, Inches(2.7), Inches(6.5), Inches(1.35),
             fill_color=BOX_FILL, corner_radius=Cm(0.15))
    add_rect(slide, MARGIN, Inches(2.7), Pt(5), Inches(1.35), fill_color=SECONDARY)
    add_text_box(slide, MARGIN + Inches(0.15), Inches(2.75), Inches(3.0), Inches(0.25),
                 "Obstacle Layer", font_name=FONT_TITLE, font_size=Pt(13),
                 color=SECONDARY, bold=True)
    add_text_box(slide, MARGIN + Inches(0.15), Inches(3.0), Inches(6.2), Inches(0.9),
                 "Rintangan dinamis dari data scan LiDAR real-time:\n"
                 "\u2022 Konversi endpoint laser ke koordinat sensor\n"
                 "\u2022 Transformasi ke frame peta\n"
                 "\u2022 Markasi occupied \u2192 cost 254 (LETHAL)\n"
                 "\u2022 Bresenham ray tracing untuk clear area antara robot dan "
                 "endpoint (set cost = 0 di sepanjang ray)",
                 font_size=Pt(10), color=BODY, line_spacing=1.25)

    # Inflation layer card
    add_rect(slide, MARGIN, Inches(4.2), Inches(6.5), Inches(0.85),
             fill_color=BOX_FILL, corner_radius=Cm(0.15))
    add_rect(slide, MARGIN, Inches(4.2), Pt(5), Inches(0.85), fill_color=SECONDARY)
    add_text_box(slide, MARGIN + Inches(0.15), Inches(4.23), Inches(3.0), Inches(0.25),
                 "Inflation Layer", font_name=FONT_TITLE, font_size=Pt(13),
                 color=SECONDARY, bold=True)
    add_text_box(slide, MARGIN + Inches(0.15), Inches(4.48), Inches(6.2), Inches(0.5),
                 "Fungsi inflasi sama dengan global costmap, namun diterapkan "
                 "pada obstacle dari scan LiDAR real-time, bukan dari peta tetap.",
                 font_size=Pt(10), color=BODY, line_spacing=1.25)

    # Right side: comparison summary
    right_x = Inches(7.5)
    add_rect(slide, right_x, Inches(1.1), Inches(5.3), Inches(1.5),
             fill_color=BOX_FILL, corner_radius=Cm(0.15))
    add_text_box(slide, right_x + Inches(0.15), Inches(1.15), Inches(5.0), Inches(0.25),
                 "Perbandingan Global vs Local", font_name=FONT_TITLE,
                 font_size=Pt(13), color=HEADING, bold=True)
    add_text_box(slide, right_x + Inches(0.15), Inches(1.45), Inches(5.0), Inches(1.0),
                 "Global Costmap:\n"
                 "\u2022 Sumber: peta tetap (static layer)\n"
                 "\u2022 Cakupan: seluruh area\n"
                 "\u2022 Statis selama navigasi\n\n"
                 "Local Costmap:\n"
                 "\u2022 Sumber: scan LiDAR real-time (obstacle layer)\n"
                 "\u2022 Cakupan: rolling window di sekitar robot\n"
                 "\u2022 Diperbarui setiap siklus kendali",
                 font_size=Pt(10), color=BODY, line_spacing=1.3)

    # Key fact box
    add_rect(slide, right_x, Inches(2.85), Inches(5.3), Inches(2.2),
             fill_color=RGBColor(0xFF, 0xF3, 0xE0), corner_radius=Cm(0.15))
    add_rect(slide, right_x, Inches(2.85), Pt(4), Inches(2.2), fill_color=ORANGE)
    add_text_box(slide, right_x + Inches(0.15), Inches(2.9), Inches(5.0), Inches(0.25),
                 "Bresenham Ray Tracing", font_name=FONT_TITLE, font_size=Pt(12),
                 color=ORANGE, bold=True)
    add_text_box(slide, right_x + Inches(0.15), Inches(3.2), Inches(5.0), Inches(1.7),
                 "Setiap sinar laser dari robot ke endpoint ditelusuri "
                 "menggunakan algoritma Bresenham 2D:\n\n"
                 "1. Inisialisasi error = \u0394x \u2212 \u0394y\n"
                 "2. Iterasi sepanjang ray dari (x\u2080, y\u2080) ke (x\u2081, y\u2081)\n"
                 "3. Setiap sel yang dilewati diberi cost = 0 (FREE_SPACE)\n"
                 "4. Endpoint diberi cost = 254 (LETHAL_OBSTACLE)\n\n"
                 "Memastikan area antara robot dan obstacle terdeteksi "
                 "sebagai ruang bebas, bukan rintangan.",
                 font_size=Pt(9), color=BODY, line_spacing=1.3)

    add_slide_number(slide, 14)


def slide_costmap_transition(prs):
    """Slide 15: Transition Costmap."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_gradient_bg(slide, GRAD_TECH_S, GRAD_TECH_E)

    add_text_box(slide, MARGIN, Inches(0.35), Inches(7), Inches(0.55),
                 "TRANSITION COSTMAP \u2014 PERPINDAHAN ANTAR REGION",
                 font_name=FONT_TITLE, font_size=Pt(26),
                 color=HEADING, bold=True)
    add_divider(slide, MARGIN, Inches(0.85), Inches(2.5))

    # Left content
    add_text_box(slide, MARGIN, Inches(1.1), Inches(6.5), Inches(0.5),
                 "Transition costmap menggabungkan data transisi dengan peta "
                 "statis. Digunakan oleh region switcher untuk mendeteksi kapan "
                 "robot memasuki perbatasan antar region dan melakukan lompatan pose.",
                 font_size=Pt(11), color=BODY, line_spacing=1.35)

    # Transition layer card
    add_rect(slide, MARGIN, Inches(1.75), Inches(6.5), Inches(1.3),
             fill_color=BOX_FILL, corner_radius=Cm(0.15))
    add_rect(slide, MARGIN, Inches(1.75), Pt(5), Inches(1.3), fill_color=ORANGE)
    add_text_box(slide, MARGIN + Inches(0.15), Inches(1.8), Inches(3.0), Inches(0.25),
                 "Transition Layer", font_name=FONT_TITLE, font_size=Pt(13),
                 color=ORANGE, bold=True)
    add_text_box(slide, MARGIN + Inches(0.15), Inches(2.1), Inches(6.2), Inches(0.8),
                 "Dua jenis nilai cost yang saling melengkapi:\n\n"
                 "T = 1   Garis transisi utama \u2014 tempat lompatan pose dieksekusi\n"
                 "E = 5   Zona buffer \u2014 robot mendeteksi bahwa ia mendekati "
                 "batas region dan mulai bersiap transisi",
                 font_size=Pt(10), color=BODY, line_spacing=1.25)

    # Region switcher flow
    add_rect(slide, MARGIN, Inches(3.2), Inches(6.5), Inches(0.6),
             fill_color=BOX_FILL, corner_radius=Cm(0.10))
    add_text_box(slide, MARGIN + Inches(0.15), Inches(3.23), Inches(6.2), Inches(0.55),
                 "Region Switcher: Deteksi E=5 \u2192 Bersiap transisi  |  "
                 "Deteksi T=1 \u2192 Eksekusi lompatan pose ke region tujuan",
                 font_size=Pt(9), color=BODY, line_spacing=1.25)

    # Static layer note
    add_rect(slide, MARGIN, Inches(3.95), Inches(6.5), Inches(0.5),
             fill_color=BOX_FILL, corner_radius=Cm(0.10))
    add_rect(slide, MARGIN, Inches(3.95), Pt(4), Inches(0.5), fill_color=SECONDARY)
    add_text_box(slide, MARGIN + Inches(0.15), Inches(3.98), Inches(6.2), Inches(0.45),
                 "Static Layer: obstacle tetap dari peta (sama dengan global costmap)",
                 font_size=Pt(9), color=SUBTLE_TEXT, line_spacing=1.25)

    # Right side: BFS expansion
    right_x = Inches(7.5)
    add_rect(slide, right_x, Inches(1.1), Inches(5.3), Inches(3.3),
             fill_color=BOX_FILL, corner_radius=Cm(0.15))
    add_text_box(slide, right_x + Inches(0.15), Inches(1.15), Inches(5.0), Inches(0.3),
                 "BFS Expansion \u2014 Zona Buffer", font_name=FONT_TITLE,
                 font_size=Pt(13), color=HEADING, bold=True)
    add_text_box(slide, right_x + Inches(0.15), Inches(1.5), Inches(5.0), Inches(2.8),
                 "Nilai E = 5 dihasilkan melalui ekspansi BFS dari sel-sel "
                 "transisi (T = 1):\n\n"
                 "1. Seed: seluruh sel dengan cost T = 1\n"
                 "2. Ekspansi: BFS 4-way hingga radius R sel\n"
                 "3. Setiap sel yang dikunjungi (bukan seed) \u2192 E = 5\n"
                 "4. Berhenti jika bertemu LETHAL_OBSTACLE (254)\n\n"
                 "Ilustrasi (R = 2):\n"
                 "  Sebelum:        Setelah:\n"
                 "  . . . . .      . E E E .\n"
                 "  . . T . .      E E T E E\n"
                 "  . T T T .  \u2192  E T T T E\n"
                 "  . . T . .      E E T E E\n"
                 "  . . . . .      . E E E .\n\n"
                 "T = 1 (garis transisi)   E = 5 (buffer deteksi)",
                 font_size=Pt(9), color=BODY, line_spacing=1.2)

    # Key fact box
    add_rect(slide, right_x, Inches(4.65), Inches(5.3), Inches(0.55),
             fill_color=RGBColor(0xFF, 0xF3, 0xE0), corner_radius=Cm(0.10))
    add_rect(slide, right_x, Inches(4.65), Pt(4), Inches(0.55), fill_color=ORANGE)
    add_text_box(slide, right_x + Inches(0.15), Inches(4.68), Inches(5.0), Inches(0.5),
                 "Transition costmap memungkinkan robot berpindah region secara "
                 "mulus tanpa kehilangan informasi konektivitas antar lantai.",
                 font_size=Pt(9), color=BODY, italic=True, line_spacing=1.25)

    add_slide_number(slide, 15)


def slide_icp(prs):
    """Slide 16: Lokalisasi ICP & Region Switcher."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_gradient_bg(slide, GRAD_TECH_S, GRAD_TECH_E)

    add_text_box(slide, MARGIN, Inches(0.35), Inches(7), Inches(0.55),
                 "LOKALISASI ICP & REGION SWITCHER", font_name=FONT_TITLE, font_size=Pt(26),
                 color=HEADING, bold=True)
    add_divider(slide, MARGIN, Inches(0.85), Inches(2.5))

    # Left: ICP pipeline
    add_text_box(slide, MARGIN, Inches(1.1), Inches(3), Inches(0.3),
                 "ICP (Iterative Closest Point)", font_name=FONT_TITLE,
                 font_size=Pt(16), color=PRIMARY, bold=True)

    icp_steps = [
        ("1", "Prekomputasi KNN", "Hash table, lookup O(1)\nk=8, r=10"),
        ("2", "Correspondence\nMatching", "Pasangkan titik\nscan → map"),
        ("3", "SVD Solver", "Estimasi rotasi &\ntranslasi optimal"),
        ("4", "Transformasi", "Perbarui pose robot\n(R, t)"),
    ]

    step_x = MARGIN
    step_y = Inches(1.55)
    step_w = Inches(2.0)
    step_h = Inches(1.4)

    for i, (num, title, desc) in enumerate(icp_steps):
        sx = step_x + i * (step_w + Inches(0.15))
        card = add_rect(slide, sx, step_y, step_w, step_h,
                        fill_color=BOX_FILL, corner_radius=Cm(0.15))
        # Number
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, sx + Inches(0.05), step_y + Inches(0.05),
            Inches(0.35), Inches(0.35)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = PRIMARY
        circle.line.fill.background()
        tf = circle.text_frame
        p = tf.paragraphs[0]
        p.text = num
        p.font.size = Pt(12)
        p.font.color.rgb = WHITE
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

        add_text_box(slide, sx + Inches(0.45), step_y + Inches(0.05),
                     step_w - Inches(0.55), Inches(0.35),
                     title, font_name=FONT_TITLE, font_size=Pt(10),
                     color=HEADING, bold=True, line_spacing=1.2)
        add_text_box(slide, sx + Inches(0.15), step_y + Inches(0.55),
                     step_w - Inches(0.3), step_h - Inches(0.65),
                     desc, font_size=Pt(9), color=BODY, line_spacing=1.35)

        # Arrow between steps
        if i < len(icp_steps) - 1:
            arr = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW,
                sx + step_w + Inches(0.01), step_y + Inches(0.55),
                Inches(0.13), Inches(0.15)
            )
            arr.fill.solid()
            arr.fill.fore_color.rgb = ORANGE
            arr.line.fill.background()

    # Parameters
    add_text_box(slide, MARGIN, Inches(3.15), Inches(11.5), Inches(0.25),
                 "Parameter: max_iter=50, k=8, r=10, t_norm=0.01m, r_norm=0.01rad, s=2",
                 font_size=Pt(10), color=BODY, italic=True)

    # Right side: Region Switcher
    add_text_box(slide, MARGIN, Inches(3.6), Inches(4), Inches(0.3),
                 "Region Switcher & Zona Transisi", font_name=FONT_TITLE,
                 font_size=Pt(16), color=ORANGE, bold=True)

    # Region switching diagram
    regions = ["Region A", "Region B", "Region C"]
    rx = MARGIN + Inches(0.5)
    ry = Inches(4.1)

    for i, rname in enumerate(regions):
        rx_pos = rx + i * Inches(3.8)
        rcard = add_rect(slide, rx_pos, ry, Inches(2.8), Inches(1.1),
                         fill_color=BOX_FILL, corner_radius=Cm(0.15))
        add_text_box(slide, rx_pos + Inches(0.1), ry + Inches(0.1),
                     Inches(2.6), Inches(0.3),
                     rname, font_name=FONT_TITLE, font_size=Pt(14),
                     color=PRIMARY, bold=True, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, rx_pos + Inches(0.1), ry + Inches(0.45),
                     Inches(2.6), Inches(0.5),
                     "Peta 2D region\nNavigasi di dalam region",
                     font_size=Pt(10), color=SUBTLE_TEXT,
                     alignment=PP_ALIGN.CENTER, line_spacing=1.3)

    # Transition arrows
    for i in range(2):
        tx_x = rx + (i + 1) * Inches(3.8) - Inches(0.8)
        tarr = slide.shapes.add_shape(
            MSO_SHAPE.RIGHT_ARROW,
            tx_x, ry + Inches(0.4), Inches(0.55), Inches(0.25)
        )
        tarr.fill.solid()
        tarr.fill.fore_color.rgb = ORANGE
        tarr.line.fill.background()

        # Transition label
        add_text_box(slide, tx_x - Inches(0.1), ry + Inches(0.7),
                     Inches(0.75), Inches(0.25),
                     "Transisi", font_size=Pt(8), color=ORANGE,
                     bold=True, alignment=PP_ALIGN.CENTER)

    # Key fact box
    key_box = add_rect(slide, MARGIN, Inches(5.5), Inches(11.5), Inches(0.6),
                       fill_color=RGBColor(0xFF, 0xF3, 0xE0), corner_radius=Cm(0.1))
    add_rect(slide, MARGIN, Inches(5.5), Pt(4), Inches(0.6), fill_color=ORANGE)
    add_text_box(slide, MARGIN + Inches(0.15), Inches(5.52), Inches(11.2), Inches(0.55),
                 "Zona transisi memerlukan radius minimal 0.26 m agar robot dapat berpindah region "
                 "dengan aman sebelum kualitas lokalisasi ICP menurun signifikan.",
                 font_size=Pt(11), color=BODY, line_spacing=1.3)

    add_slide_number(slide, 16)


def slide_path_planning(prs):
    """Slide 17: Path Planning (A* Global Planner)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_gradient_bg(slide, GRAD_TECH_S, GRAD_TECH_E)

    add_text_box(slide, MARGIN, Inches(0.35), Inches(7), Inches(0.55),
                 "PATH PLANNING — A* GLOBAL PLANNER", font_name=FONT_TITLE, font_size=Pt(26),
                 color=HEADING, bold=True)
    add_divider(slide, MARGIN, Inches(0.85), Inches(2.5))

    # --- A* Algorithm ---
    add_text_box(slide, MARGIN, Inches(1.1), Inches(3), Inches(0.3),
                 "Algoritma A*", font_name=FONT_TITLE, font_size=Pt(18),
                 color=PRIMARY, bold=True)

    astar_content = (
        "• f(n) = g(n) + h(n)\n"
        "    g(n) = biaya dari start ke node n\n"
        "    h(n) = estimasi heuristik dari n ke goal\n\n"
        "• Quadratic potential calculation untuk\n"
        "  menghitung cost gradient pada grid\n\n"
        "• Open list / Closed list untuk eksplorasi\n"
        "  node secara optimal"
    )
    add_rect(slide, MARGIN - Inches(0.10), Inches(1.45),
             Inches(5.5), Inches(2.1),
             fill_color=BOX_FILL, corner_radius=Cm(0.10))
    add_text_box(slide, MARGIN, Inches(1.5), Inches(5.3), Inches(2.0),
                 astar_content, font_size=Pt(12), color=BODY, line_spacing=1.4)

    # --- Global Costmap Integration ---
    add_text_box(slide, MARGIN, Inches(3.65), Inches(5), Inches(0.3),
                 "Integrasi dengan Global Costmap", font_name=FONT_TITLE,
                 font_size=Pt(14), color=PRIMARY, bold=True)

    global_content = (
        "Global planner menggunakan static layer dari global\n"
        "costmap untuk membangun cost gradient pada occupancy\n"
        "grid. Setiap sel memiliki biaya navigasi berdasarkan\n"
        "jarak ke obstacle terdekat dan Cost Scaling Factor.\n\n"
        "Inflasi obstacle menghasilkan gradien biaya halus yang\n"
        "mendorong path menjauhi rintangan. Fungsi biaya:\n"
        "  C(d) = 253 · e^(−α · (d − r))  untuk d ≤ R"
    )
    add_rect(slide, MARGIN - Inches(0.10), Inches(3.95),
             Inches(5.5), Inches(2.7),
             fill_color=BOX_FILL, corner_radius=Cm(0.10))
    add_text_box(slide, MARGIN, Inches(4.05), Inches(5.3), Inches(2.5),
                 global_content, font_size=Pt(11), color=BODY, line_spacing=1.35)

    # --- Right side: Path Interceptor + Multi-Region diagram ---
    right_x = Inches(7.0)

    add_text_box(slide, right_x, Inches(1.1), Inches(5.5), Inches(0.3),
                 "Path Interceptor (Multi-Region)", font_name=FONT_TITLE,
                 font_size=Pt(18), color=ORANGE, bold=True)

    # Diagram: 3 region boxes connected by transition arrows
    regions = ["Region A", "Region B", "Region C"]
    rx = right_x + Inches(0.3)
    ry = Inches(1.55)
    rw = Inches(2.0)
    rh = Inches(0.75)
    rgap = Inches(0.3)

    for i, rname in enumerate(regions):
        # Vertical stacking with transition in between
        rcard = add_rect(slide, rx, ry + i * (rh + rgap),
                         rw, rh, fill_color=BOX_FILL, corner_radius=Cm(0.10))
        add_rect(slide, rx, ry + i * (rh + rgap), rw, Pt(3), fill_color=PRIMARY)
        add_text_box(slide, rx + Inches(0.05), ry + i * (rh + rgap) + Inches(0.03),
                     rw - Inches(0.1), Inches(0.3),
                     rname, font_name=FONT_TITLE, font_size=Pt(13),
                     color=PRIMARY, bold=True, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, rx + Inches(0.05),
                     ry + i * (rh + rgap) + Inches(0.35),
                     rw - Inches(0.1), Inches(0.35),
                     "A* planning pada\ncostmap region",
                     font_size=Pt(9), color=SUBTLE_TEXT,
                     alignment=PP_ALIGN.CENTER, line_spacing=1.25)

        if i < 2:
            # Arrow between regions
            arr = slide.shapes.add_shape(
                MSO_SHAPE.DOWN_ARROW,
                rx + rw / 2 - Inches(0.15),
                ry + (i + 0.5) * (rh + rgap) - Inches(0.05),
                Inches(0.3), Inches(0.35)
            )
            arr.fill.solid()
            arr.fill.fore_color.rgb = ORANGE
            arr.line.fill.background()

    # Right side bottom: Path Interceptor explanation
    add_rect(slide, right_x, Inches(4.55), Inches(5.5), Inches(2.1),
             fill_color=BOX_FILL, corner_radius=Cm(0.10))
    add_rect(slide, right_x, Inches(4.55), Pt(4), Inches(2.1), fill_color=ORANGE)
    add_text_box(slide, right_x + Inches(0.15), Inches(4.65), Inches(5.2), Inches(1.9),
                 "Path interceptor memecah rute global multi-region\n"
                 "menjadi segmen per region yang di-plan secara\n"
                 "independen oleh A*:\n\n"
                 "  Region A  →  Titik Transisi A-B  →\n"
                 "  Region B  →  Titik Transisi B-C  →\n"
                 "  Region C  →  Goal\n\n"
                 "Setiap segmen menggunakan costmap region masing-\n"
                 "masing. Transition costmap menyediakan informasi\n"
                 "konektor untuk perpindahan mulus antar region.",
                 font_size=Pt(11), color=BODY, line_spacing=1.35)

    add_slide_number(slide, 17)


def slide_dwa(prs):
    """Slide 18: Path Tracking (DWA Local Planner)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_gradient_bg(slide, GRAD_TECH_S, GRAD_TECH_E)

    add_text_box(slide, MARGIN, Inches(0.35), Inches(7), Inches(0.55),
                 "PATH TRACKING — DWA LOCAL PLANNER", font_name=FONT_TITLE, font_size=Pt(26),
                 color=HEADING, bold=True)
    add_divider(slide, MARGIN, Inches(0.85), Inches(2.5))

    # Top: 5-step flow
    steps = [
        ("1. Velocity\nSampling", "6 vx × 6 vy × 30 ω\n= 1080 trajectory"),
        ("2. Kinematic\nSimulation", "Simulasi robot\nselama Δt"),
        ("3. Trajectory\nScoring", "7 kriteria\npenilaian"),
        ("4. Best\nSelection", "Pilih trajectory\ndengan biaya terendah"),
        ("5. Send to\nRobot", "Eksekusi gerakan\npada hexapod"),
    ]

    sx = MARGIN
    sy = Inches(1.15)
    sw = Inches(2.1)
    sh = Inches(1.45)

    for i, (title, desc) in enumerate(steps):
        scard = add_rect(slide, sx, sy, sw, sh,
                         fill_color=BOX_FILL, corner_radius=Cm(0.15))
        add_rect(slide, sx, sy, sw, Pt(4), fill_color=PRIMARY)
        add_text_box(slide, sx + Inches(0.1), sy + Inches(0.1),
                     sw - Inches(0.2), Inches(0.5),
                     title, font_name=FONT_TITLE, font_size=Pt(10),
                     color=HEADING, bold=True, alignment=PP_ALIGN.CENTER,
                     line_spacing=1.2)
        add_text_box(slide, sx + Inches(0.1), sy + Inches(0.65),
                     sw - Inches(0.2), Inches(0.7),
                     desc, font_size=Pt(9), color=SUBTLE_TEXT,
                     alignment=PP_ALIGN.CENTER, line_spacing=1.35)

        if i < len(steps) - 1:
            arr = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW,
                sx + sw + Inches(0.02), sy + Inches(0.55),
                Inches(0.22), Inches(0.2)
            )
            arr.fill.solid()
            arr.fill.fore_color.rgb = ORANGE
            arr.line.fill.background()

        sx += sw + Inches(0.26)

    # Bottom: 7 scoring criteria
    add_text_box(slide, MARGIN, Inches(2.8), Inches(11.5), Inches(0.3),
                 "7 Kriteria Scoring DWA", font_name=FONT_TITLE,
                 font_size=Pt(16), color=HEADING, bold=True)

    criteria = [
        ("C_osc", "Oscillation\npenalty"),
        ("C_obs", "Obstacle\ndistance"),
        ("C_gf", "Goal front\nalignment"),
        ("C_align", "Path\nalignment"),
        ("C_path", "Distance to\nglobal path"),
        ("C_goal", "Distance\nto goal"),
        ("C_twirl", "Rotation\npenalty"),
    ]

    weights = ["w_osc", "w_obs", "w_gf", "w_align", "w_path", "w_goal", "w_twirl"]
    colors_crit = [PRIMARY, SECONDARY, SECONDARY, SECONDARY, SECONDARY, SECONDARY, ORANGE]

    cx = MARGIN
    cy = Inches(3.25)
    cw = Inches(1.55)
    ch = Inches(1.1)

    for i, ((c_name, c_desc), w_name, color) in enumerate(zip(criteria, weights, colors_crit)):
        ccard = add_rect(slide, cx, cy, cw, ch,
                         fill_color=BOX_FILL, corner_radius=Cm(0.1))
        add_text_box(slide, cx + Inches(0.05), cy + Inches(0.05),
                     Inches(1.1), Inches(0.2),
                     c_name, font_size=Pt(10), color=color, bold=True)
        add_text_box(slide, cx + Inches(0.05), cy + Inches(0.28),
                     cw - Inches(0.1), Inches(0.45),
                     c_desc, font_size=Pt(8), color=SUBTLE_TEXT,
                     line_spacing=1.3)
        add_text_box(slide, cx + Inches(0.05), cy + Inches(0.78),
                     cw - Inches(0.1), Inches(0.2),
                     w_name, font_size=Pt(7), color=BODY)
        cx += cw + Inches(0.1)

    # Formula
    add_rect(slide, MARGIN - Inches(0.10), Inches(4.48),
             Inches(11.7), Inches(0.45),
             fill_color=BOX_FILL, corner_radius=Cm(0.10))
    add_text_box(slide, MARGIN, Inches(4.55), Inches(11.5), Inches(0.35),
                 "C_total = C_osc + Σ(wi · Ci)  →  Minimalkan total cost",
                 font_size=Pt(12), color=BODY, bold=True, alignment=PP_ALIGN.CENTER)

    # Key feature
    key_box = add_rect(slide, MARGIN, Inches(5.15), Inches(11.5), Inches(0.55),
                       fill_color=RGBColor(0xFF, 0xF3, 0xE0), corner_radius=Cm(0.1))
    add_text_box(slide, MARGIN + Inches(0.15), Inches(5.2), Inches(11.2), Inches(0.45),
                 "Mampu menghindari obstacle dinamis yang tidak ada di global path, "
                 "tercermin dari peningkatan max cross track error pada kondisi gempa.",
                 font_size=Pt(11), color=BODY, line_spacing=1.3)

    add_slide_number(slide, 18)


def slide_arena(prs):
    """Slide 19: Arena Pengujian & Skenario."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_gradient_bg(slide, GRAD_RESULTS_S, GRAD_RESULTS_E)

    add_text_box(slide, MARGIN, Inches(0.35), Inches(6), Inches(0.55),
                 "ARENA PENGUJIAN & SKENARIO", font_name=FONT_TITLE, font_size=Pt(28),
                 color=HEADING, bold=True)
    add_divider(slide, MARGIN, Inches(0.85), Inches(2.5))

    base_dir = os.path.dirname(os.path.abspath(__file__))
    img_h = Inches(2.2)
    img_w = int(img_h * (1847.0 / 1051.0))
    gap_x = Inches(0.45)
    gap_y = Inches(0.35)
    total_w = img_w * 2 + gap_x
    start_x = (SLIDE_W - total_w) / 2

    arena_imgs = [
        ("sim_RA.png", "Region A \u2014 Lorong & Ramp", start_x, Inches(1.2)),
        ("sim_RA_rough.png", "Region A \u2014 Rough Floor", start_x + img_w + gap_x, Inches(1.2)),
        ("sim_RB.png", "Region B \u2014 Lorong Sempit", start_x, Inches(1.2) + img_h + gap_y),
        ("sim_RC2.png", "Region C \u2014 Area Terbuka", start_x + img_w + gap_x, Inches(1.2) + img_h + gap_y),
    ]

    for fname, label, ax, ay in arena_imgs:
        img_path = os.path.join(base_dir, "gambar", fname)
        add_text_box(slide, ax, ay - Inches(0.22),
                     img_w, Inches(0.20),
                     label, font_size=Pt(9), color=PRIMARY, bold=True,
                     alignment=PP_ALIGN.CENTER)
        slide.shapes.add_picture(img_path, ax, ay, img_w, img_h)

    add_text_box(slide, MARGIN, Inches(6.15), Inches(11.5), Inches(0.35),
                 "Platform: Mujoco Simulator  \u00b7  Robot: Hexapod (6 kaki)  \u00b7  "
                 "Sensor: LiDAR 2D  \u00b7  Peta: Known Map",
                 font_size=Pt(11), color=BODY, bold=False)

    add_slide_number(slide, 19)


def slide_results_icp(prs):
    """Slide 20: Hasil Pengujian — ICP Localization."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_gradient_bg(slide, GRAD_RESULTS_S, GRAD_RESULTS_E)

    add_text_box(slide, MARGIN, Inches(0.35), Inches(7), Inches(0.55),
                 "HASIL PENGUJIAN: ICP LOCALIZATION", font_name=FONT_TITLE, font_size=Pt(26),
                 color=HEADING, bold=True)
    add_divider(slide, MARGIN, Inches(0.85), Inches(2.5))

    base_dir = os.path.dirname(os.path.abspath(__file__))
    img_h = Inches(3.0)
    img_w = int(img_h * (1500.0 / 1200.0))
    gap = Inches(0.5)
    total_w = img_w * 2 + gap
    start_x = (SLIDE_W - total_w) / 2
    img_top = Inches(1.15)

    icp_imgs = [
        ("icp_errors_RA.png", "Region A — Normal Floor", start_x),
        ("icp_errors_RA_rough.png", "Region A — Rough Floor", start_x + img_w + gap),
    ]

    for fname, label, ix in icp_imgs:
        img_path = os.path.join(base_dir, "gambar", fname)
        add_text_box(slide, ix, img_top - Inches(0.25),
                     img_w, Inches(0.22),
                     label, font_size=Pt(10), color=PRIMARY, bold=True,
                     alignment=PP_ALIGN.CENTER)
        slide.shapes.add_picture(img_path, ix, img_top, img_w, img_h)

    # Key finding
    add_rect(slide, MARGIN, Inches(4.50), Inches(11.5), Inches(0.55),
             fill_color=RGBColor(0xFF, 0xF3, 0xE0), corner_radius=Cm(0.1))
    add_rect(slide, MARGIN, Inches(4.50), Pt(4), Inches(0.55), fill_color=ORANGE)
    add_text_box(slide, MARGIN + Inches(0.15), Inches(4.52), Inches(11.2), Inches(0.50),
                 "ICP akurat pada lantai normal. Rough floor menyebabkan osilasi "
                 "signifikan pada sumbu x dan yaw akibat variasi vertikal scan LiDAR.",
                 font_size=Pt(11), color=BODY, line_spacing=1.3)

    add_slide_number(slide, 20)


def _make_path_result_slide(prs, region_label, map_png, csf_csv, stats_csv, slide_num,
                             key_finding, key_color=ORANGE):
    """Generic path planning result slide for a single region."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_gradient_bg(slide, GRAD_RESULTS_S, GRAD_RESULTS_E)

    add_text_box(slide, MARGIN, Inches(0.3), Inches(10), Inches(0.5),
                 f"HASIL PATH PLANNING \u2014 REGION {region_label}", font_name=FONT_TITLE,
                 font_size=Pt(24), color=HEADING, bold=True)
    add_divider(slide, MARGIN, Inches(0.75), Inches(2.5))

    base_dir = os.path.dirname(os.path.abspath(__file__))

    # --- Left side: map image ---
    img_path = os.path.join(base_dir, "gambar", map_png)
    img_h = Inches(4.5)
    img_w = img_h  # square-ish fallback
    try:
        from PIL import Image
        im = Image.open(img_path)
        img_w = int(img_h * (im.width / im.height))
    except Exception:
        pass
    img_left = MARGIN
    img_top = Inches(1.05)
    slide.shapes.add_picture(img_path, img_left, img_top, img_w, img_h)

    # --- Right side: tables + key finding ---
    right_x = img_left + img_w + Inches(0.4)
    tbl_left = right_x
    tbl_top = Inches(1.05)

    csf_widths = [Inches(0.6), Inches(1.0), Inches(1.0), Inches(1.0), Inches(1.0)]
    _add_csv_table(slide, tbl_left, tbl_top, csf_widths, csf_csv,
                   title="Tabel CSF \u2014 Panjang, Waktu, Clearance, Waypoint")

    stats_top = tbl_top + Inches(1.7)
    stats_widths = [Inches(0.6), Inches(1.0), Inches(1.0), Inches(1.0)]
    _add_csv_table(slide, tbl_left, stats_top, stats_widths, stats_csv,
                   title="Tabel Deviasi \u2014 Mean/Max Cross Track & SD")

    key_y = stats_top + Inches(1.8)
    add_rect(slide, right_x, key_y, Inches(4.6), Inches(1.0),
             fill_color=RGBColor(0xFF, 0xF3, 0xE0), corner_radius=Cm(0.1))
    add_rect(slide, right_x, key_y, Pt(4), Inches(1.0), fill_color=key_color)
    add_text_box(slide, right_x + Inches(0.15), key_y + Inches(0.05),
                 Inches(4.3), Inches(0.9),
                 key_finding, font_size=Pt(10), color=BODY, line_spacing=1.35)

    add_slide_number(slide, slide_num)
    return slide


def slide_results_path_RA(prs):
    """Slide 21: Hasil Path Planning — Region A."""
    _make_path_result_slide(prs, "A", "map_with_paths_RA.png",
                            "path_CSF_RA.csv", "path_stats_RA.csv", 21,
                            "CSF 100 = path terpendek (0.97 m) & paling efisien.\n"
                            "CSF 10/20/50 = path lebih panjang, di tengah lorong.\n"
                            "SD CSF 100 paling rendah (0.002 m) -> konsistensi tinggi.",
                            key_color=RGBColor(0x2E, 0x7D, 0x32))


def slide_results_path_RB(prs):
    """Slide 22: Hasil Path Planning — Region B."""
    _make_path_result_slide(prs, "B", "map_with_paths_RB.png",
                            "path_CSF_RB.csv", "path_stats_RB.csv", 22,
                            "CSF tidak signifikan: layout lorong sempit -> semua\n"
                            "CSF menghasilkan path hampir identik (~0.59 m).\n"
                            "SD sangat kecil (< 0.004 m) untuk semua CSF.",
                            key_color=RGBColor(0x2E, 0x7D, 0x32))


def slide_results_path_RC(prs):
    """Slide 23: Hasil Path Planning — Region C."""
    _make_path_result_slide(prs, "C", "map_with_paths_RC.png",
                            "path_CSF_RC.csv", "path_stats_RC.csv", 23,
                            "CSF 100 = path terpendek (2.52 m).\n"
                            "CSF 10 lebih aman di tengah lorong, lebih panjang.\n"
                            "Trade-off: path pendek vs deviasi ground truth.",
                            key_color=RGBColor(0x2E, 0x7D, 0x32))

def slide_results_navigation(prs):
    """Slide 24: Hasil Pengujian — Navigasi Penuh."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_gradient_bg(slide, GRAD_RESULTS_S, GRAD_RESULTS_E)

    add_text_box(slide, MARGIN, Inches(0.3), Inches(10), Inches(0.5),
                 "HASIL PENGUJIAN: NAVIGASI PENUH", font_name=FONT_TITLE, font_size=Pt(24),
                 color=HEADING, bold=True)
    add_divider(slide, MARGIN, Inches(0.72), Inches(2.5))

    # --- Left side: DWA tracking tables ---
    tbl_left = MARGIN
    tbl_top = Inches(1.05)

    track_widths = [Inches(0.9), Inches(1.15), Inches(1.55), Inches(1.55)]
    tbl1, h1, _ = _add_csv_table(slide, tbl_left, tbl_top, track_widths,
                                  "path_track.csv",
                                  title="Cross Track Error — Kondisi Normal")

    tbl2_top = tbl_top + h1 + Inches(0.25)
    _add_csv_table(slide, tbl_left, tbl2_top, track_widths,
                   "path_track_rough.csv",
                   title="Cross Track Error — Kondisi Rough Floor")

    # --- Right side: DWA summary + navigation highlight ---
    right_x = Inches(6.7)

    # DWA summary cards
    dwa_results = [
        ("Kondisi Normal", "Mean & max cross track error rendah.\nDWA mengikuti global path dengan baik.",
         RGBColor(0x2E, 0x7D, 0x32), "\u2713"),
        ("Kondisi Gempa", "Max cross track error meningkat signifikan.\nObstacle tidak diketahui di global path \u2192\nDWA menyimpang untuk menghindari tabrakan.",
         ORANGE, "\u26A0"),
    ]

    dry = Inches(1.0)
    for title, desc, color, status in dwa_results:
        drcard = add_rect(slide, right_x, dry, Inches(5.9), Inches(1.0),
                          fill_color=BOX_FILL, corner_radius=Cm(0.1))
        add_rect(slide, right_x, dry, Pt(4), Inches(1.0), fill_color=color)
        add_text_box(slide, right_x + Inches(0.12), dry + Inches(0.04),
                     Inches(2.0), Inches(0.22),
                     f"{status}  {title}", font_name=FONT_TITLE, font_size=Pt(12),
                     color=color, bold=True)
        add_text_box(slide, right_x + Inches(0.12), dry + Inches(0.28),
                     Inches(5.6), Inches(0.65),
                     desc, font_size=Pt(9), color=BODY, line_spacing=1.35)
        dry += Inches(1.1)

    # Navigation highlight box
    hl_top = Inches(3.25)
    result_box = add_rect(slide, right_x, hl_top, Inches(5.9), Inches(1.6),
                          fill_color=RGBColor(0xE3, 0xF0, 0xFD), corner_radius=Cm(0.2))
    add_rect(slide, right_x, hl_top, Inches(5.9), Pt(4), fill_color=PRIMARY)

    add_text_box(slide, right_x, hl_top + Inches(0.05), Inches(5.9), Inches(0.28),
                 "NAVIGASI PENUH (Region A \u2192 C)",
                 font_name=FONT_TITLE, font_size=Pt(14), color=PRIMARY, bold=True,
                 alignment=PP_ALIGN.CENTER)

    add_text_box(slide, right_x, hl_top + Inches(0.35), Inches(5.9), Inches(0.5),
                 "127.26",
                 font_name=FONT_TITLE, font_size=Pt(42), color=PRIMARY, bold=True,
                 alignment=PP_ALIGN.CENTER)

    add_text_box(slide, right_x, hl_top + Inches(0.85), Inches(5.9), Inches(0.2),
                 "DETIK",
                 font_name=FONT_TITLE, font_size=Pt(14), color=PRIMARY, bold=True,
                 alignment=PP_ALIGN.CENTER)

    add_text_box(slide, right_x, hl_top + Inches(1.1), Inches(5.9), Inches(0.35),
                 "Robot berhasil menavigasi Region A \u2192 C secara otonom penuh.",
                 font_size=Pt(10), color=BODY, alignment=PP_ALIGN.CENTER)

    # Component integration box
    int_top = Inches(5.0)
    int_box = add_rect(slide, right_x, int_top, Inches(5.9), Inches(1.35),
                       fill_color=BOX_FILL, corner_radius=Cm(0.15))
    add_text_box(slide, right_x + Inches(0.15), int_top + Inches(0.05),
                 Inches(5.6), Inches(0.2),
                 "Semua Komponen Bekerja Terintegrasi",
                 font_name=FONT_TITLE, font_size=Pt(11), color=HEADING, bold=True,
                 alignment=PP_ALIGN.CENTER)

    integrations = [
        "\u2713  ICP Localization \u2192 estimasi pose real-time",
        "\u2713  Costmap Transisi \u2192 perpindahan antar region",
        "\u2713  A* Global Planner \u2192 rute multi-region",
        "\u2713  DWA Local Planner \u2192 trajectory + obstacle avoidance",
    ]

    inty = int_top + Inches(0.3)
    for item in integrations:
        add_text_box(slide, right_x + Inches(0.2), inty, Inches(5.5), Inches(0.2),
                     item, font_size=Pt(9), color=BODY, line_spacing=1.1)
        inty += Inches(0.22)

    add_slide_number(slide, 24)


def slide_conclusion(prs):
    """Slide 25: Kesimpulan."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_gradient_bg(slide, GRAD_CLOSING_S, GRAD_CLOSING_E)

    add_text_box(slide, MARGIN, Inches(0.35), Inches(3.5), Inches(0.55),
                 "KESIMPULAN", font_name=FONT_TITLE, font_size=Pt(30),
                 color=HEADING, bold=True)
    add_divider(slide, MARGIN, Inches(0.85), Inches(1.8))

    conclusions = [
        ("1", "Layered costmap berhasil merepresentasikan lingkungan\nmulti-lantai sebagai kumpulan peta 2D region", SECONDARY),
        ("2", "ICP localization akurat untuk navigasi, dengan peningkatan\nerror pada lantai tidak rata dan region kompleks", SECONDARY),
        ("3", "Cost Scaling Factor 100 menghasilkan jalur terpendek\ndengan tetap menghindari obstacle", ORANGE),
        ("4", "Costmap transisi berhasil menggabungkan region pada\nradius transisi minimal 0.26 meter", SECONDARY),
        ("5", "DWA Planner mampu mengakomodasi obstacle tidak\ndiketahui dan menghasilkan trajectory aman", ORANGE),
    ]

    card_w = Inches(3.55)
    card_h = Inches(2.2)
    positions = [
        (MARGIN, Inches(1.15)),
        (MARGIN + card_w + Inches(0.25), Inches(1.15)),
        (MARGIN + 2 * (card_w + Inches(0.25)), Inches(1.15)),
        (MARGIN + Inches(1.9), Inches(3.6)),
        (MARGIN + card_w + Inches(0.25) + Inches(1.9), Inches(3.6)),
    ]

    for (x, y), (num, desc, color) in zip(positions, conclusions):
        card = add_rect(slide, x, y, card_w, card_h,
                        fill_color=BOX_FILL, corner_radius=Cm(0.15))
        # Big number background
        num_bg = add_rect(slide, x, y, card_w, Inches(0.6), fill_color=color,
                          corner_radius=Cm(0.15))
        add_text_box(slide, x, y + Inches(0.05), card_w, Inches(0.5),
                     num, font_name=FONT_TITLE, font_size=Pt(26),
                     color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
        add_text_box(slide, x + Inches(0.2), y + Inches(0.8),
                     card_w - Inches(0.4), Inches(1.2),
                     desc, font_size=Pt(12), color=BODY, line_spacing=1.45)

    add_slide_number(slide, 25)


def slide_future_work(prs):
    """Slide 26: Saran / Future Work."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_gradient_bg(slide, GRAD_CLOSING_S, GRAD_CLOSING_E)

    add_text_box(slide, MARGIN, Inches(0.35), Inches(4), Inches(0.55),
                 "SARAN PENGEMBANGAN", font_name=FONT_TITLE, font_size=Pt(30),
                 color=HEADING, bold=True)
    add_divider(slide, MARGIN, Inches(0.85), Inches(2.0))

    suggestions = [
        ("1", "Prekomputasi KNN per Region",
         "Parameter radius dan jumlah nearest neighbors\ndapat disetel independen per region untuk\nakurasi ICP yang lebih optimal.",
         PRIMARY),
        ("2", "Fusion ICP + IMU",
         "Kompensasi lonjakan error ICP di zona transisi\ndengan data IMU saat fitur LiDAR terbatas.",
         SECONDARY),
        ("3", "Integrasi Depth Camera",
         "Keterbatasan LiDAR 2D dalam mendeteksi\nobstacle vertikal dapat dilengkapi dengan\ndata visual dari depth sensor.",
         SECONDARY),
        ("4", "Pengujian pada Robot Fisik",
         "Validasi sistem di lingkungan nyata setelah\nterbukti berhasil di lingkungan simulasi Mujoco.",
         ORANGE),
    ]

    card_w = Inches(5.6)
    card_h = Inches(2.1)
    positions_sug = [
        (MARGIN, Inches(1.15)),
        (MARGIN + card_w + Inches(0.3), Inches(1.15)),
        (MARGIN, Inches(3.5)),
        (MARGIN + card_w + Inches(0.3), Inches(3.5)),
    ]

    for (x, y), (num, title, desc, color) in zip(positions_sug, suggestions):
        card = add_rect(slide, x, y, card_w, card_h,
                        fill_color=BOX_FILL, corner_radius=Cm(0.15))
        add_rect(slide, x, y, card_w, Pt(4), fill_color=color)
        add_text_box(slide, x + Inches(0.2), y + Inches(0.1),
                     Inches(0.4), Inches(0.35),
                     num, font_name=FONT_TITLE, font_size=Pt(20),
                     color=color, bold=True)
        add_text_box(slide, x + Inches(0.6), y + Inches(0.1),
                     card_w - Inches(0.8), Inches(0.35),
                     title, font_name=FONT_TITLE, font_size=Pt(13),
                     color=HEADING, bold=True)
        add_text_box(slide, x + Inches(0.2), y + Inches(0.65),
                     card_w - Inches(0.4), Inches(1.3),
                     desc, font_size=Pt(11), color=BODY, line_spacing=1.5)

    add_slide_number(slide, 26)


def slide_thanks(prs):
    """Slide 27: Penutup / Thank You."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    # Decorative line
    add_rect(slide, SLIDE_W / 2 - Inches(1.5), Inches(1.2), Inches(3.0), Pt(3), fill_color=ORANGE)

    add_text_box(slide, Inches(0), Inches(1.5), SLIDE_W, Inches(1.0),
                 "TERIMA KASIH", font_name=FONT_TITLE, font_size=Pt(42),
                 color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(0), Inches(2.8), SLIDE_W, Inches(0.5),
                 "Bagas Surya Wirawan  ·  NRP. 5022221026",
                 font_name=FONT_TITLE, font_size=Pt(18),
                 color=RGBColor(0xBB, 0xD5, 0xEC),
                 alignment=PP_ALIGN.CENTER)

    add_text_box(slide, Inches(0), Inches(3.5), SLIDE_W, Inches(1.5),
                 "Dosen Pembimbing:\n"
                 "Dr. Ir. Djoko Purwanto, M.Eng.\n"
                 "Fajar Budiman, S.T., M.Sc.",
                 font_name=FONT_BODY, font_size=Pt(14),
                 color=RGBColor(0x99, 0xBB, 0xDD),
                 alignment=PP_ALIGN.CENTER, line_spacing=1.4)

    add_text_box(slide, Inches(0), Inches(5.3), SLIDE_W, Inches(1.0),
                 "Departemen Teknik Elektro\n"
                 "Fakultas Teknologi Elektro dan Informatika Cerdas\n"
                 "Institut Teknologi Sepuluh Nopember\n"
                 "Surabaya · 2026",
                 font_name=FONT_BODY, font_size=Pt(12),
                 color=RGBColor(0x88, 0xAA, 0xCC),
                 alignment=PP_ALIGN.CENTER, line_spacing=1.4)

    # Bottom decorative bars
    for i, color in enumerate([PRIMARY, SECONDARY, ORANGE]):
        bar_w = Inches(3.5)
        add_rect(slide, SLIDE_W / 2 - Inches(5.2) + i * (bar_w + Inches(0.1)),
                 SLIDE_H - Inches(0.35), bar_w, Pt(2), fill_color=color)


# =============================================================================
# MAIN
# =============================================================================

def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # Build all slides
    slide_title(prs)
    slide_outline(prs)
    slide_background(prs)
    slide_problem(prs)
    slide_constraints(prs)
    slide_objectives(prs)
    slide_literature(prs)
    slide_prior_research(prs)
    slide_system_overview(prs)
    slide_nav_flow(prs)
    slide_map_processing(prs)
    slide_costmap_overview(prs)
    slide_costmap_global(prs)
    slide_costmap_local(prs)
    slide_costmap_transition(prs)
    slide_icp(prs)
    slide_path_planning(prs)
    slide_dwa(prs)
    slide_arena(prs)
    slide_results_icp(prs)
    slide_results_path_RA(prs)
    slide_results_path_RB(prs)
    slide_results_path_RC(prs)
    slide_results_navigation(prs)
    slide_conclusion(prs)
    slide_future_work(prs)
    slide_thanks(prs)

    output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "Presentasi_Sidang_TA.pptx")
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
